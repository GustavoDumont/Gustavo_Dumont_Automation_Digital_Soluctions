"""
I LOVE Cmt 2.1 — Suite de ferramentas PDF
Concremat Engenharia | Desenvolvido por Gustavo Dumont

Correções v2.1 aplicadas:
 [1]  Início não divide mais a tela ao clicar múltiplas vezes
 [2]  Organizar/Reordenar: lógica original preservada
 [3]  Dividir PDF: lógica original preservada
 [8]  PDF→Word: log de erros visível na tela
 [9]  PDF→Excel: calibragem de tolerância ajustável
 [10] Editar PDF: posicionamento interativo por clique no preview
 [14] Marca D'água: posicionamento interativo por clique no preview
 [15] Assinatura Digital: reconstruída via PIL (sem retângulo preto)
 [16] Descomprimir PDF: nova função adicionada
 [17] Renomeador: não divide tela + Lista Auxiliar 2

Instalação:
  pip install customtkinter pypdf pillow PyMuPDF reportlab
              pytesseract rapidfuzz python-docx docx2pdf
              pdfplumber openpyxl python-pptx pdf2image
"""

# ── Imports padrão ────────────────────────────────────────────────────────────
import os, sys, re, io, shutil, tempfile, threading, time
import tkinter as tk
import importlib.util
from pathlib import Path
from typing import List, Tuple
from io import BytesIO

import customtkinter as ctk
from tkinter import filedialog, messagebox, simpledialog, colorchooser
from pypdf import PdfReader, PdfWriter
from PIL import Image, ImageTk, ImageDraw
import fitz  # PyMuPDF

from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.utils import ImageReader
from reportlab.lib import colors as rl_colors

# ── Imports opcionais ─────────────────────────────────────────────────────────
try:
    from docx2pdf import convert as docx2pdf_convert
except Exception:
    docx2pdf_convert = None

try:
    if getattr(sys, "frozen", False):
        BASE_DIR = sys._MEIPASS
    else:
        BASE_DIR = os.path.dirname(os.path.abspath(__file__))
except Exception:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

TESSERACT_PORTABLE = os.path.join(BASE_DIR, "tesseract", "tesseract.exe")
TESSDATA_PORTABLE  = os.path.join(BASE_DIR, "tesseract", "tessdata")

try:
    import pytesseract
    if os.path.exists(TESSERACT_PORTABLE):
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_PORTABLE
        os.environ["TESSDATA_PREFIX"] = TESSDATA_PORTABLE
    HAVE_TESSERACT = True
except Exception:
    HAVE_TESSERACT = False

try:
    from rapidfuzz import fuzz as rfuzz
    USE_RAPIDFUZZ = True
except Exception:
    import difflib
    USE_RAPIDFUZZ = False

try:
    import pdfplumber
    HAVE_PDFPLUMBER = True
except Exception:
    HAVE_PDFPLUMBER = False

try:
    import openpyxl
    from openpyxl.styles import Font as XLFont, PatternFill
    HAVE_OPENPYXL = True
except Exception:
    HAVE_OPENPYXL = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAVE_PPTX = True
except Exception:
    HAVE_PPTX = False

try:
    from docx import Document as DocxDocument
    from docx.shared import Inches as DocxInches
    HAVE_DOCX = True
except Exception:
    HAVE_DOCX = False

if importlib.util.find_spec("docx") is None:
    try:
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-docx"])
        from docx import Document as DocxDocument
        from docx.shared import Inches as DocxInches
        HAVE_DOCX = True
    except Exception:
        pass

# ── Tema ──────────────────────────────────────────────────────────────────────
ctk.set_appearance_mode("Dark")
ctk.set_default_color_theme("blue")

APP_NAME    = "I LOVE Cmt"
APP_VERSION = "2.1"
BRAND_BLUE  = "#1B4F8A"
BRAND_BLUE2 = "#2563EB"
ACCENT      = "#3B82F6"
SIDEBAR_BG  = "#111827"
MAIN_BG     = "#1E1E2E"
CARD_BG     = "#2A2A3C"
TEXT_DIM    = "#9CA3AF"

# ── Helpers globais ───────────────────────────────────────────────────────────

def run_in_thread(fn, *args, **kwargs):
    t = threading.Thread(target=fn, args=args, kwargs=kwargs, daemon=True)
    t.start()
    return t

def unique_path(path: Path) -> Path:
    if not path.exists():
        return path
    stem, suffix = path.stem, path.suffix
    i = 1
    while True:
        p = path.parent / f"{stem}({i}){suffix}"
        if not p.exists():
            return p
        i += 1

# ── Funções de PDF ────────────────────────────────────────────────────────────

def merge_pdfs(files, out_path):
    writer = PdfWriter()
    for f in files:
        reader = PdfReader(str(f))
        for page in reader.pages:
            writer.add_page(page)
    with open(str(out_path), "wb") as fo:
        writer.write(fo)

def images_to_pdfs(image_paths, out_dir, single_pdf=False,
                    single_name="images_merged.pdf"):
    outputs = []
    pil_images = []
    for img_path in image_paths:
        img = Image.open(img_path)
        if img.mode in ("RGBA", "LA"):
            bg = Image.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[3])
            img = bg
        else:
            img = img.convert("RGB")
        if single_pdf:
            pil_images.append(img)
        else:
            name = Path(img_path).stem + ".pdf"
            op = Path(out_dir) / name
            img.save(op, "PDF", resolution=100.0)
            outputs.append(str(op))
    if single_pdf and pil_images:
        first, rest = pil_images[0], pil_images[1:]
        op = Path(out_dir) / single_name
        first.save(op, "PDF", save_all=True, append_images=rest)
        outputs.append(str(op))
    return outputs

def compress_pdf_fitz(input_path, output_path, image_dpi=72, jpeg_quality=60):
    orig_size = os.path.getsize(input_path)
    with fitz.open(input_path) as doc:
        for page in doc:
            for img_ref in page.get_images(full=True):
                xref = img_ref[0]
                try:
                    base_image = doc.extract_image(xref)
                    img_bytes  = base_image["image"]
                    img_ext    = base_image["ext"]
                    pil_img    = Image.open(BytesIO(img_bytes))
                    w, h = pil_img.size
                    scale = min(image_dpi / 150, 1.0)
                    nw, nh = max(1, int(w * scale)), max(1, int(h * scale))
                    pil_img = pil_img.resize((nw, nh), Image.LANCZOS)
                    buf = BytesIO()
                    fmt = "JPEG" if img_ext.lower() in ("jpg", "jpeg") else "PNG"
                    if fmt == "JPEG":
                        pil_img = pil_img.convert("RGB")
                        pil_img.save(buf, format="JPEG",
                                     quality=jpeg_quality, optimize=True)
                    else:
                        pil_img.save(buf, format="PNG", optimize=True)
                    buf.seek(0)
                    doc.update_stream(xref, buf.read())
                except Exception:
                    pass
        doc.save(str(output_path), garbage=4, deflate=True, clean=True)
    new_size = os.path.getsize(str(output_path))
    return int((1 - new_size / orig_size) * 100) if orig_size else 0

def decompress_pdf(input_path, output_path):
    """Remove compressão de streams do PDF."""
    with fitz.open(input_path) as doc:
        doc.save(str(output_path), garbage=0, deflate=False,
                 clean=False, expand=255)
    orig_size = os.path.getsize(input_path)
    new_size  = os.path.getsize(str(output_path))
    return new_size, orig_size

def overlay_watermark(pdf_path, watermark_pdf, out_path):
    with open(str(pdf_path), "rb") as fp, open(str(watermark_pdf), "rb") as fw:
        reader    = PdfReader(fp)
        watermark = PdfReader(fw)
        wp        = watermark.pages[0]
        writer    = PdfWriter()
        for p in reader.pages:
            p.merge_page(wp)
            writer.add_page(p)
        with open(str(out_path), "wb") as fo:
            writer.write(fo)

def apply_watermark(pdf_path, watermark_image=None, watermark_text=None,
                    output_path=None, scale_factor=1.0, alpha=0.5,
                    pos_x_pct=0.5, pos_y_pct=0.5):
    if not output_path:
        raise ValueError("output_path é obrigatório")
    tmp_pdf = Path(tempfile.gettempdir()) / f"__tmp_wm_{os.getpid()}.pdf"
    c   = rl_canvas.Canvas(str(tmp_pdf), pagesize=letter)
    w, h = letter
    if watermark_image:
        img = Image.open(watermark_image).convert("RGBA")
        nw  = int(img.width  * scale_factor)
        nh  = int(img.height * scale_factor)
        img = img.resize((max(1, nw), max(1, nh)), Image.LANCZOS)
        al  = img.split()[3].point(lambda p: int(p * alpha))
        img.putalpha(al)
        buf = BytesIO(); img.save(buf, "PNG"); buf.seek(0)
        c.drawImage(ImageReader(buf),
                    w * pos_x_pct - nw / 2, h * pos_y_pct - nh / 2,
                    width=nw, height=nh, mask="auto")
    if watermark_text:
        fs = int(96 * scale_factor)
        c.saveState()
        c.setFont("Helvetica-Bold", fs)
        c.setFillColorRGB(0, 0, 0.55)
        c.setFillAlpha(alpha)
        tw = c.stringWidth(watermark_text, "Helvetica-Bold", fs)
        c.drawString(w * pos_x_pct - tw / 2, h * pos_y_pct, watermark_text)
        c.restoreState()
    c.save()
    overlay_watermark(str(pdf_path), str(tmp_pdf), str(output_path))
    tmp_pdf.unlink(missing_ok=True)


# ═══════════════════════════════════════════════════════════════════════════════
# CLASSE PRINCIPAL
# ═══════════════════════════════════════════════════════════════════════════════

class App(ctk.CTk):

    # ──────────────────────────────────────────────────────────────────────────
    # INIT
    # ──────────────────────────────────────────────────────────────────────────

    def __init__(self):
        super().__init__()
        self.title(f"{APP_NAME} {APP_VERSION}")
        self.geometry("1200x720")
        self.minsize(1000, 620)
        self.configure(fg_color=MAIN_BG)

        self._active_btn   = None
        self.current_frame = None

        self.grid_columnconfigure(0, weight=0)
        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)
        self.grid_rowconfigure(1, weight=0)

        self._build_sidebar()
        self._build_statusbar()

        # área de conteúdo principal
        self.main = ctk.CTkFrame(self, corner_radius=0, fg_color=MAIN_BG)
        self.main.grid(row=0, column=1, sticky="nswe")
        self.main.grid_rowconfigure(0, weight=1)
        self.main.grid_columnconfigure(0, weight=1)

        self.open_home()

    # ──────────────────────────────────────────────────────────────────────────
    # SIDEBAR
    # ──────────────────────────────────────────────────────────────────────────

    def _build_sidebar(self):
        self.sidebar = ctk.CTkFrame(self, width=230, corner_radius=0,
                                     fg_color=SIDEBAR_BG)
        self.sidebar.grid(row=0, column=0, rowspan=2, sticky="nswe")
        self.sidebar.grid_propagate(False)
        self.sidebar.grid_columnconfigure(0, weight=1)

        logo = ctk.CTkFrame(self.sidebar, fg_color="transparent")
        logo.grid(row=0, column=0, pady=(18, 4), padx=12, sticky="ew")
        ctk.CTkLabel(logo, text="I LOVE Cmt",
                     font=ctk.CTkFont("Segoe UI", 20, "bold"),
                     text_color="#FFFFFF").pack()
        ctk.CTkLabel(logo, text=f"v{APP_VERSION}  •  Concremat",
                     font=ctk.CTkFont("Segoe UI", 10),
                     text_color=TEXT_DIM).pack()

        ctk.CTkFrame(self.sidebar, height=1, fg_color="#2D3748").grid(
            row=1, column=0, sticky="ew", padx=12, pady=(4, 8))

        scroll = ctk.CTkScrollableFrame(self.sidebar, fg_color="transparent",
                                         scrollbar_button_color=SIDEBAR_BG)
        scroll.grid(row=2, column=0, sticky="nsew")
        self.sidebar.grid_rowconfigure(2, weight=1)
        scroll.grid_columnconfigure(0, weight=1)

        row = [0]

        def sep(label):
            ctk.CTkLabel(scroll, text=f"  {label}",
                         font=ctk.CTkFont("Segoe UI", 10, "bold"),
                         text_color=TEXT_DIM, anchor="w").grid(
                row=row[0], column=0, sticky="w", padx=14, pady=(12, 2))
            row[0] += 1

        def btn(icon, label, cmd):
            b = ctk.CTkButton(
                scroll,
                text=f"  {icon}  {label}",
                anchor="w", height=36, corner_radius=8,
                fg_color="transparent", hover_color="#1E293B",
                text_color="#CBD5E1",
                font=ctk.CTkFont("Segoe UI", 13))
            b.configure(command=lambda c=cmd, bt=b: self._nav(c, bt))
            b.grid(row=row[0], column=0, sticky="ew", padx=8, pady=1)
            row[0] += 1
            return b

        btn("🏠", "Início",                  self.open_home)

        sep("ORGANIZAR")
        btn("📎", "Juntar PDFs",             self.open_tab_merge)
        btn("🗂️", "Organizar / Reordenar",   self.open_tab_organize)
        btn("✂️", "Dividir PDF",              self.open_tab_split_pdf)
        btn("🔄", "Girar Páginas",            self.open_tab_rotate)

        sep("CONVERTER")
        btn("🖼️",  "Imagens → PDF",           self.open_tab_images)
        btn("📸",  "PDF → Imagens",            self.open_pdf_to_images)
        btn("📝",  "Word → PDF",               self.open_tab_word)
        btn("📄",  "PDF → Word",               self.open_tab_pdf_to_word)
        btn("📊",  "PDF → Excel",              self.open_tab_pdf_to_excel)
        btn("🖥️",  "PDF → PowerPoint",         self.open_tab_pdf_to_pptx)

        sep("PROTEGER")
        btn("🔒",  "Proteger com Senha",       self.open_tab_protect)
        btn("🔓",  "Remover Senha",             self.open_tab_remove_password)
        btn("💧",  "Marca D'água",              self.open_tab_watermark)

        sep("EDITAR")
        btn("✏️",  "Editar PDF",               self.open_tab_edit_pdf)
        btn("✍️",  "Assinatura Digital",        self.open_tab_signature)

        sep("COMPRESSÃO")
        btn("🗜️",  "Comprimir PDF",             self.open_tab_compress)
        btn("📦",  "Descomprimir PDF",           self.open_tab_decompress)

        sep("RENOMEAR")
        btn("🏷️",  "Renomeador de PDFs",        self.open_tab_renamer)

        ctk.CTkLabel(self.sidebar,
                     text="© Concremat Engenharia\nDesenv. por Gustavo Dumont",
                     font=ctk.CTkFont("Segoe UI", 9),
                     text_color=TEXT_DIM, justify="center").grid(
            row=3, column=0, pady=(8, 14), padx=8)

    def _nav(self, cmd, btn_widget):
        if self._active_btn:
            try:
                self._active_btn.configure(fg_color="transparent",
                                            text_color="#CBD5E1")
            except Exception:
                pass
        self._active_btn = btn_widget
        btn_widget.configure(fg_color=BRAND_BLUE, text_color="#FFFFFF")
        cmd()

    # ──────────────────────────────────────────────────────────────────────────
    # BARRA DE STATUS
    # ──────────────────────────────────────────────────────────────────────────

    def _build_statusbar(self):
        self.statusbar = ctk.CTkFrame(self, height=28, corner_radius=0,
                                       fg_color="#0F172A")
        self.statusbar.grid(row=1, column=1, sticky="ew")
        self.statusbar.grid_propagate(False)
        self.status_label = ctk.CTkLabel(
            self.statusbar, text="Pronto.",
            font=ctk.CTkFont("Segoe UI", 11),
            text_color=TEXT_DIM, anchor="w")
        self.status_label.pack(side="left", padx=12)
        self.progress_bar = ctk.CTkProgressBar(self.statusbar, width=160, height=10)
        self.progress_bar.pack(side="right", padx=12, pady=8)
        self.progress_bar.set(0)

    def set_status(self, msg: str, progress: float = None):
        def _u():
            self.status_label.configure(text=msg)
            if progress is not None:
                self.progress_bar.set(progress)
        self.after(0, _u)

    # ──────────────────────────────────────────────────────────────────────────
    # HELPERS DE UI
    # ──────────────────────────────────────────────────────────────────────────

    def clear_main(self):
        """Destrói o frame de conteúdo atual e limpa filhos residuais."""
        if self.current_frame is not None:
            try:
                self.current_frame.destroy()
            except Exception:
                pass
            self.current_frame = None
        for w in list(self.main.winfo_children()):
            try:
                w.destroy()
            except Exception:
                pass

    def _page(self):
        f = ctk.CTkFrame(self.main, fg_color=MAIN_BG)
        f.pack(fill="both", expand=True)
        return f

    def _page_title(self, parent, icon, title, subtitle=""):
        hdr = ctk.CTkFrame(parent, fg_color="transparent")
        hdr.pack(fill="x", padx=28, pady=(22, 6))
        ctk.CTkLabel(hdr, text=f"{icon}  {title}",
                     font=ctk.CTkFont("Segoe UI", 22, "bold"),
                     anchor="w").pack(anchor="w")
        if subtitle:
            ctk.CTkLabel(hdr, text=subtitle,
                         font=ctk.CTkFont("Segoe UI", 13),
                         text_color=TEXT_DIM, anchor="w").pack(anchor="w")
        ctk.CTkFrame(parent, height=1, fg_color="#2D3748").pack(
            fill="x", padx=28, pady=(6, 0))

    def _action_btn(self, parent, text, command, width=200, color=BRAND_BLUE2):
        return ctk.CTkButton(
            parent, text=text, command=command,
            width=width, height=44, corner_radius=10,
            fg_color=color, hover_color=BRAND_BLUE,
            font=ctk.CTkFont("Segoe UI", 14, "bold"))

    def _card(self, parent, **kw):
        return ctk.CTkFrame(parent, corner_radius=12, fg_color=CARD_BG, **kw)

    def _file_listbox(self, parent, height=150):
        return ctk.CTkTextbox(
            parent, height=height,
            font=ctk.CTkFont("Consolas", 12),
            fg_color="#12121E", text_color="#A0AEC0",
            corner_radius=8)

    # ══════════════════════════════════════════════════════════════════════════
    # 🏠 HOME  — FIX: não divide a tela ao clicar múltiplas vezes
    # ══════════════════════════════════════════════════════════════════════════

    def open_home(self):
        self.clear_main()

        outer = ctk.CTkFrame(self.main, fg_color=MAIN_BG)
        outer.pack(fill="both", expand=True)
        self.current_frame = outer          # salva UMA VEZ

        canvas  = tk.Canvas(outer, bg=MAIN_BG, highlightthickness=0)
        vscroll = tk.Scrollbar(outer, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vscroll.set)
        vscroll.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        inner = tk.Frame(canvas, bg=MAIN_BG)
        win_id = canvas.create_window((0, 0), window=inner, anchor="nw")

        def _resize(e=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfig(win_id, width=canvas.winfo_width())

        inner.bind("<Configure>", _resize)
        canvas.bind("<Configure>", _resize)
        canvas.bind_all("<MouseWheel>",
            lambda e: canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        self.set_status("Bem-vindo ao I LOVE Cmt 2.1", 0)

        # Hero
        hero = tk.Frame(inner, bg=BRAND_BLUE)
        hero.pack(fill="x", padx=28, pady=(28, 20))
        tk.Label(hero, text=f"I LOVE Cmt  {APP_VERSION}",
                 font=("Segoe UI", 28, "bold"),
                 fg="#FFFFFF", bg=BRAND_BLUE).pack(pady=(24, 4))
        tk.Label(hero,
                 text="Suite completa de ferramentas · Concremat Engenharia",
                 font=("Segoe UI", 13), fg="#BFDBFE", bg=BRAND_BLUE).pack(
            pady=(0, 24))

        # Atalhos
        tk.Label(inner, text="  ⚡  Acesso Rápido",
                 font=("Segoe UI", 14, "bold"),
                 fg="#E2E8F0", bg=MAIN_BG).pack(anchor="w", padx=28, pady=(4, 8))

        grid_f = tk.Frame(inner, bg=MAIN_BG)
        grid_f.pack(fill="x", padx=28)

        shortcuts = [
            ("📎", "Juntar PDFs",         self.open_tab_merge),
            ("🗂️", "Organizar Páginas",   self.open_tab_organize),
            ("💧", "Marca D'água",         self.open_tab_watermark),
            ("🗜️", "Comprimir PDF",        self.open_tab_compress),
            ("📦", "Descomprimir PDF",     self.open_tab_decompress),
            ("🔒", "Proteger com Senha",   self.open_tab_protect),
            ("✏️", "Editar PDF",           self.open_tab_edit_pdf),
            ("📊", "PDF → Excel",          self.open_tab_pdf_to_excel),
            ("🏷️", "Renomeador",           self.open_tab_renamer),
        ]
        for i, (ic, lbl, cmd) in enumerate(shortcuts):
            tk.Button(grid_f, text=f"{ic}\n{lbl}",
                       command=cmd, relief="flat", cursor="hand2",
                       font=("Segoe UI", 12), fg="#E2E8F0", bg=CARD_BG,
                       activebackground="#334155", width=14, height=3
                       ).grid(row=i // 3, column=i % 3,
                               padx=6, pady=6, sticky="ew")
        for c in range(3):
            grid_f.grid_columnconfigure(c, weight=1)

        # Novidades
        tk.Label(inner, text="  ✨  Novidades na versão 2.1",
                 font=("Segoe UI", 14, "bold"),
                 fg="#E2E8F0", bg=MAIN_BG).pack(
            anchor="w", padx=28, pady=(24, 8))

        news_items = [
            ("📦", "Descomprimir PDF",       "Remove compressão para facilitar edição"),
            ("📊", "PDF → Excel calibrado",  "Tolerância ajustável para tabelas limpas"),
            ("✏️", "Editar PDF interativo",   "Clique no preview para posicionar conteúdo"),
            ("💧", "Marca D'água interativa", "Posicione a marca clicando no preview"),
            ("✍️", "Assinatura Digital",       "Desenho vetorial inserido sem fundo preto"),
            ("🏷️", "Renomeador — 3 listas",   "Nova Lista Auxiliar 2 para cruzamento duplo"),
            ("📄", "PDF→Word com log",         "Veja erros página a página na tela"),
            ("🔄", "Girar páginas",            "Selecione páginas ou angulos específicos"),
        ]
        news_g = tk.Frame(inner, bg=MAIN_BG)
        news_g.pack(fill="x", padx=28, pady=(0, 32))
        news_g.grid_columnconfigure((0, 1), weight=1, uniform="col")
        for i, (ic, ttl, desc) in enumerate(news_items):
            c = tk.Frame(news_g, bg=CARD_BG, padx=14, pady=10)
            c.grid(row=i // 2, column=i % 2, padx=6, pady=6, sticky="ew")
            tk.Label(c, text=f"{ic} {ttl}",
                     font=("Segoe UI", 12, "bold"),
                     fg="#E2E8F0", bg=CARD_BG, anchor="w").pack(anchor="w")
            tk.Label(c, text=desc,
                     font=("Segoe UI", 11), fg=TEXT_DIM,
                     bg=CARD_BG, anchor="w", wraplength=300).pack(anchor="w")

    # ══════════════════════════════════════════════════════════════════════════
    # 📎 JUNTAR PDFs
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_merge(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "📎", "Juntar PDFs",
                          "Combine vários PDFs em um único arquivo.")

        card = self._card(frame)
        card.pack(fill="both", expand=True, padx=28, pady=14)

        info = (
            "📄 Instruções:\n\n"
            "1️⃣  Clique em 'Adicionar PDFs' para selecionar os arquivos.\n"
            "2️⃣  A ordem seguirá a seleção feita no explorador.\n"
            "3️⃣  Clique em 'Juntar e Salvar' para gerar o PDF combinado."
        )
        ctk.CTkLabel(card, text=info, justify="left", wraplength=680,
                     font=ctk.CTkFont("Segoe UI", 13),
                     text_color=TEXT_DIM).pack(anchor="w", padx=16, pady=(14, 8))

        self.merge_listbox = self._file_listbox(card, height=180)
        self.merge_listbox.pack(fill="both", expand=True, padx=16, pady=4)
        self.merge_files: List[str] = []

        btns = ctk.CTkFrame(card, fg_color="transparent")
        btns.pack(fill="x", padx=16, pady=(4, 8))
        self._action_btn(btns, "➕ Adicionar PDFs",
                          self._merge_add, width=180).pack(side="left", padx=(0, 8))
        self._action_btn(btns, "✖ Limpar lista",
                          lambda: (self.merge_files.clear() or
                                   self.merge_listbox.delete("1.0", "end")),
                          width=150, color="#7F1D1D").pack(side="left")
        self._action_btn(btns, "🧩 Juntar e Salvar",
                          self._merge_run, width=180).pack(side="right")

        self.merge_progress = ctk.CTkProgressBar(card)
        self.merge_progress.pack(fill="x", padx=16, pady=(0, 14))
        self.merge_progress.set(0)

    def _merge_add(self):
        files = filedialog.askopenfilenames(title="Selecionar PDFs",
                                             filetypes=[("PDF", "*.pdf")])
        if not files:
            return
        self.merge_files.extend(files)
        self.merge_listbox.delete("1.0", "end")
        for f in self.merge_files:
            self.merge_listbox.insert("end", f"{Path(f).name}\n")

    def _merge_run(self):
        if len(self.merge_files) < 2:
            messagebox.showwarning("Aviso", "Selecione ao menos 2 PDFs.")
            return
        out = filedialog.asksaveasfilename(
            title="Salvar PDF combinado",
            defaultextension=".pdf",
            filetypes=[("PDF", "*.pdf")],
            initialfile="PDFs_Unidos.pdf")
        if not out:
            return

        def worker():
            self.set_status("Juntando PDFs...", 0)
            try:
                merge_pdfs(self.merge_files, out)
                self.merge_progress.set(1.0)
                self.set_status(f"PDF salvo: {Path(out).name}", 1)
                self.after(0, lambda: messagebox.showinfo(
                    "Concluído", f"PDFs unidos com sucesso!\n{out}"))
            except Exception as e:
                self.after(0, lambda: messagebox.showerror("Erro", str(e)))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 🗂️ ORGANIZAR — lógica original completa
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_organize(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "🗂️", "Organizar / Reordenar Páginas",
                          "Reordene páginas por arrastar e soltar.")

        card = self._card(frame)
        card.pack(fill="x", padx=28, pady=14)
        ctk.CTkLabel(card, text=(
            "1. Selecione o PDF para organizar.\n"
            "2. Clique em uma miniatura para visualizar à direita.\n"
            "3. Arraste e solte as miniaturas para reordenar.\n"
            "4. Clique em 'Salvar PDF Organizado'."),
            justify="left", wraplength=700,
            font=ctk.CTkFont("Segoe UI", 13),
            text_color=TEXT_DIM).pack(padx=16, pady=(14, 8))

        self._action_btn(card, "🗂️ Abrir PDF para Organizar",
                          self.start_organize_ui, width=280).pack(pady=(4, 18))

    def start_organize_ui(self):
        """Lógica original de organização com drag-and-drop."""
        from PIL import ImageDraw, ImageFont
        from pypdf import PdfReader, PdfWriter

        file_path = filedialog.askopenfilename(
            title="Selecionar PDF para organizar",
            filetypes=[("Arquivos PDF", "*.pdf")])
        if not file_path:
            return

        org_window = ctk.CTkToplevel(self)
        org_window.title(f"Organizar páginas — {Path(file_path).name}")
        org_window.geometry("1280x800")

        frame_main = ctk.CTkFrame(org_window)
        frame_main.pack(fill="both", expand=True, padx=10, pady=10)

        left_frame = ctk.CTkFrame(frame_main)
        left_frame.pack(side="left", fill="both", expand=True, padx=(0, 10))

        right_frame = ctk.CTkFrame(frame_main, width=520)
        right_frame.pack(side="right", fill="y")

        ctk.CTkLabel(right_frame, text="Pré-visualização",
                     font=ctk.CTkFont(size=16, weight="bold")).pack(pady=10)
        preview_canvas = tk.Canvas(right_frame, bg="#1e1e1e",
                                    width=480, height=650, highlightthickness=0)
        preview_canvas.pack(padx=10, pady=(0, 10))

        _preview_state = {"current_page": None, "render_thread": None,
                           "cancel_flag": False}
        _drag_state    = {"start_x": 0, "start_y": 0, "widget": None,
                           "index": None, "moved": False}
        thumb_widgets  = []

        # ---- helpers ----
        def convert_page_silent(path, page_number, target_width=None, dpi=150):
            try:
                with fitz.open(path) as doc:
                    p   = doc.load_page(page_number - 1)
                    mat = fitz.Matrix(dpi / 72.0, dpi / 72.0)
                    pix = p.get_pixmap(matrix=mat, alpha=False)
                    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                if target_width:
                    w, h = img.size
                    img  = img.resize((int(target_width),
                                       int(h * target_width / w)), Image.LANCZOS)
                return img
            except Exception as e:
                print(f"[convert err] {e}")
                return None

        def add_page_number_overlay(pil_img, page_number):
            draw = ImageDraw.Draw(pil_img)
            fs   = max(14, int(pil_img.width * 0.06))
            try:
                font = ImageFont.truetype("arial.ttf", fs)
            except Exception:
                font = ImageFont.load_default()
            text = str(page_number)
            try:
                bb = draw.textbbox((0, 0), text, font=font)
                tw, th = bb[2]-bb[0], bb[3]-bb[1]
            except Exception:
                tw, th = fs, fs
            draw.rectangle([(0, 0), (tw+12, th+12)], fill=(0, 0, 0, 150))
            draw.text((6, 6), text, fill="white", font=font)
            return pil_img

        def safe_update_preview(path, page):
            if (_preview_state["render_thread"] and
                    _preview_state["render_thread"].is_alive()):
                _preview_state["cancel_flag"] = True
            _preview_state["cancel_flag"] = False
            _preview_state["current_page"] = page
            preview_canvas.delete("all")
            preview_canvas.create_text(
                preview_canvas.winfo_reqwidth()//2,
                preview_canvas.winfo_reqheight()//2,
                text=f"Carregando página {page}...",
                fill="gray", font=("Arial", 12, "italic"))
            preview_canvas.update_idletasks()

            def _render():
                try:
                    doc = fitz.open(path)
                    if _preview_state["cancel_flag"]:
                        doc.close(); return
                    pobj = doc.load_page(page - 1)
                    pix  = pobj.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
                    doc.close()
                    if _preview_state["cancel_flag"]: return
                    pil = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    cw  = preview_canvas.winfo_width()  or 480
                    ch  = preview_canvas.winfo_height() or 650
                    sc  = min(cw/pil.width, ch/pil.height)
                    pil = pil.resize((max(1,int(pil.width*sc)),
                                      max(1,int(pil.height*sc))), Image.LANCZOS)
                    tkimg = ImageTk.PhotoImage(pil)
                    def _u():
                        if _preview_state["cancel_flag"]: return
                        preview_canvas.delete("all")
                        preview_canvas.create_image(cw//2, ch//2,
                                                     image=tkimg, anchor="center")
                        preview_canvas.image = tkimg
                    preview_canvas.after(0, _u)
                except Exception as e:
                    print(f"[preview err] {e}")
            t = threading.Thread(target=_render, daemon=True)
            _preview_state["render_thread"] = t
            t.start()

        # ---- drag & drop ----
        def on_thumb_press(event, page):
            w = event.widget
            try:
                idx = thumb_widgets.index(w)
            except ValueError:
                idx = None
            _drag_state.update({"start_x": event.x_root, "start_y": event.y_root,
                                  "widget": w, "index": idx, "moved": False})
            w.config(relief="ridge")

        def on_thumb_motion(event):
            if not _drag_state["widget"]: return
            if (abs(event.x_root - _drag_state["start_x"]) > 8 or
                    abs(event.y_root - _drag_state["start_y"]) > 8):
                _drag_state["moved"] = True

        def on_thumb_release(event):
            w = _drag_state["widget"]
            if not w: return
            w.config(relief="solid")
            if not _drag_state["moved"]:
                page = getattr(w, "page_num", None)
                if page is not None:
                    safe_update_preview(file_path, page)
            else:
                try:
                    x  = event.x_root - inner_frame.winfo_rootx()
                    y  = event.y_root - inner_frame.winfo_rooty()
                    dc = max(0, int(x // 280))
                    dr = max(0, int(y // 360))
                    di = min(dr*3+dc, len(thumb_widgets)-1)
                    si = _drag_state["index"]
                    if si is None:
                        _drag_state.update({"widget":None,"index":None,"moved":False})
                        return
                    item = thumb_widgets.pop(si)
                    thumb_widgets.insert(di, item)
                    for idx, ww in enumerate(thumb_widgets):
                        ww.grid(row=idx//3, column=idx%3, padx=8, pady=8)
                except Exception as e:
                    print(f"[drag release err] {e}")
            _drag_state.update({"widget": None, "index": None, "moved": False})

        def create_thumbnail_button(parent, pil_image, page_num):
            tkimg = ImageTk.PhotoImage(pil_image)
            lbl   = tk.Label(parent, image=tkimg, bg="#2b2b2b",
                              bd=2, relief="solid", cursor="hand2")
            lbl.image    = tkimg
            lbl.page_num = page_num
            lbl.grid(row=(page_num-1)//3, column=(page_num-1)%3, padx=8, pady=8)
            lbl.bind("<ButtonPress-1>",   lambda e, p=page_num: on_thumb_press(e, p))
            lbl.bind("<B1-Motion>",        on_thumb_motion)
            lbl.bind("<ButtonRelease-1>",  on_thumb_release)
            lbl.bind("<Enter>",  lambda e: lbl.config(bd=3, relief="ridge"))
            lbl.bind("<Leave>",  lambda e: lbl.config(bd=2, relief="solid"))
            return lbl

        # ---- scroll area ----
        cvs      = tk.Canvas(left_frame, bg="#1e1e1e", highlightthickness=0)
        sbar     = ctk.CTkScrollbar(left_frame, orientation="vertical",
                                     command=cvs.yview)
        cvs.configure(yscrollcommand=sbar.set)
        sbar.pack(side="right", fill="y")
        cvs.pack(side="left", fill="both", expand=True, padx=(0, 8))
        inner_frame = tk.Frame(cvs, bg="#2b2b2b")
        cvs.create_window((0, 0), window=inner_frame, anchor="nw")
        inner_frame.bind("<Configure>",
                          lambda e: cvs.configure(scrollregion=cvs.bbox("all")))

        def on_mw(e): cvs.yview_scroll(int(-1*(e.delta/120)), "units")
        cvs.bind("<Enter>",  lambda e: org_window.bind_all("<MouseWheel>", on_mw))
        cvs.bind("<Leave>",  lambda e: org_window.unbind_all("<MouseWheel>"))

        reader  = PdfReader(file_path)
        n_pages = len(reader.pages)

        def gen_thumbnails():
            for i in range(n_pages):
                pidx = i + 1
                pil  = convert_page_silent(file_path, pidx,
                                            target_width=260, dpi=90)
                if pil is None: continue
                pil = add_page_number_overlay(pil, pidx)
                def mk(pi=pidx, p=pil):
                    lbl = create_thumbnail_button(inner_frame, p, pi)
                    if lbl:
                        thumb_widgets.append(lbl)
                        cvs.configure(scrollregion=cvs.bbox("all"))
                org_window.after(0, mk)
            org_window.after(200, lambda: cvs.configure(
                scrollregion=cvs.bbox("all")))

        threading.Thread(target=gen_thumbnails, daemon=True).start()

        def save_new_pdf():
            try:
                order = [w.page_num for w in thumb_widgets]
                out   = filedialog.asksaveasfilename(
                    defaultextension=".pdf",
                    filetypes=[("Arquivos PDF", "*.pdf")],
                    title="Salvar PDF organizado")
                if not out: return
                writer = PdfWriter()
                for idx in order:
                    writer.add_page(reader.pages[idx - 1])
                with open(out, "wb") as f:
                    writer.write(f)
                messagebox.showinfo("Sucesso", f"PDF salvo em:\n{out}")
                org_window.destroy()
            except Exception as e:
                messagebox.showerror("Erro", str(e))

        ctk.CTkButton(org_window, text="💾 Salvar PDF Organizado",
                       command=save_new_pdf, fg_color=BRAND_BLUE2,
                       font=ctk.CTkFont(size=14, weight="bold"),
                       height=44).pack(pady=12, padx=20, fill="x")

    # ══════════════════════════════════════════════════════════════════════════
    # ✂️ DIVIDIR PDF — lógica original preservada
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_split_pdf(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "✂️", "Dividir PDF",
                          "Divida um PDF em partes por intervalos, "
                          "páginas individuais ou a cada X páginas.")
        card = self._card(frame)
        card.pack(fill="x", padx=28, pady=14)
        ctk.CTkLabel(card, text=(
            "Selecione um PDF para abrir a tela visual de divisão.\n"
            "Você poderá excluir páginas, dividir por intervalos "
            "ou a cada X páginas."),
            wraplength=700, font=ctk.CTkFont("Segoe UI", 13),
            text_color=TEXT_DIM).pack(padx=16, pady=(14, 8))
        self._action_btn(card, "✂️ Selecionar PDF para Dividir",
                          self.run_split_pdf, width=280).pack(pady=(4, 18))

    def run_split_pdf(self):
        fp = filedialog.askopenfilename(title="Selecione um PDF",
                                         filetypes=[("PDF", "*.pdf")])
        if fp:
            self.start_visual_pdf_ui(file_path=fp, mode="split")

    def start_visual_pdf_ui(self, file_path: str, mode: str = "organize"):
        """Lógica original completa de visualização/divisão de PDF."""
        win = ctk.CTkToplevel(self)
        win.geometry("1300x800")
        win.title("🧩 Organizar PDF" if mode == "organize" else "✂️ Dividir PDF")
        win.configure(fg_color="#1E1E1E")

        def on_close():
            try:
                if hasattr(self, "_pdf_doc") and self._pdf_doc:
                    self._pdf_doc.close()
            except Exception:
                pass
            win.destroy()
        win.protocol("WM_DELETE_WINDOW", on_close)

        # top bar
        top_bar = ctk.CTkFrame(win, height=60)
        top_bar.pack(side="top", fill="x", padx=10, pady=(10, 5))

        if mode == "split":
            split_mode = tk.StringVar(value="interval")
            ctk.CTkRadioButton(top_bar, text="Intervalos",
                variable=split_mode, value="interval").pack(side="left", padx=6)
            ctk.CTkRadioButton(top_bar, text="Páginas individuais",
                variable=split_mode, value="single").pack(side="left", padx=6)
            ctk.CTkRadioButton(top_bar, text="A cada X páginas",
                variable=split_mode, value="every").pack(side="left", padx=6)
            interval_entry = ctk.CTkEntry(top_bar,
                placeholder_text="Ex: 1-3,5,7-9", width=200)
            interval_entry.pack(side="left", padx=10)
            every_entry = ctk.CTkEntry(top_bar,
                placeholder_text="X páginas", width=120)
            every_entry.pack(side="left", padx=6)
            ctk.CTkButton(top_bar, text="🗑 Excluir Página", width=160,
                           fg_color="#7F1D1D", hover_color="#991B1B",
                           command=lambda: delete_selected_page()
                           ).pack(side="right", padx=6)
            ctk.CTkButton(top_bar, text="✂️ Dividir PDF", width=160,
                           fg_color=BRAND_BLUE2, hover_color=BRAND_BLUE,
                           command=lambda: self.execute_split_visual(
                               split_mode.get(),
                               interval_entry.get(),
                               every_entry.get())
                           ).pack(side="right", padx=6)
        else:
            ctk.CTkButton(top_bar, text="💾 Salvar PDF", width=160,
                           fg_color="#16A34A", hover_color="#166534",
                           command=lambda: self.save_organized_pdf(win)
                           ).pack(side="right", padx=6)

        # layout
        left = tk.Frame(win, bg="#1E1E1E", width=680)
        left.pack(side="left", fill="y", padx=(10, 4), pady=10)
        left.pack_propagate(False)

        cvs = tk.Canvas(left, bg="#1E1E1E", highlightthickness=0)
        cvs.pack(side="left", fill="both", expand=True)
        sb  = tk.Scrollbar(left, orient="vertical", command=cvs.yview)
        sb.pack(side="right", fill="y", padx=(6, 0))
        cvs.configure(yscrollcommand=sb.set)

        inner = tk.Frame(cvs, bg="#1E1E1E")
        drag_indicator = tk.Frame(inner, height=4, bg=ACCENT)
        drag_indicator.place_forget()
        cvs.create_window((0, 0), window=inner, anchor="nw")
        inner.bind("<Configure>", lambda _: cvs.configure(
            scrollregion=cvs.bbox("all")))

        right_fr = ctk.CTkFrame(win, fg_color="#111111")
        right_fr.pack(side="left", fill="both", expand=True,
                       padx=(4, 10), pady=10)
        preview = tk.Canvas(right_fr, bg="#111111", highlightthickness=0)
        preview.pack(fill="both", expand=True, padx=6, pady=6)

        cvs.bind_all("<MouseWheel>",
            lambda e: cvs.yview_scroll(-1*int(e.delta/120), "units"))

        self._organize_pdf_path = file_path
        self._thumb_frames  = []
        self._thumb_images  = []
        self._orig_page_idx = []
        self._selected_thumb = None
        self._pdf_doc = fitz.open(self._organize_pdf_path)

        def pil_from_page(page_idx):
            with fitz.open(self._organize_pdf_path) as d:
                p   = d.load_page(page_idx)
                pix = p.get_pixmap(dpi=90)
            return Image.open(BytesIO(pix.tobytes("png"))).convert("RGB")

        def show_preview(idx):
            def _r():
                img = pil_from_page(idx)
                w   = preview.winfo_width()  or 600
                h   = preview.winfo_height() or 800
                r   = min(w/img.width, h/img.height)
                img = img.resize((int(img.width*r), int(img.height*r)), Image.LANCZOS)
                tkimg = ImageTk.PhotoImage(img)
                def _u():
                    preview.delete("all")
                    preview.create_image((w-img.width)//2, (h-img.height)//2,
                                          anchor="nw", image=tkimg)
                    preview.image = tkimg
                win.after(0, _u)
            run_in_thread(_r)

        def delete_selected_page():
            if not self._selected_thumb:
                messagebox.showwarning("Aviso", "Nenhuma página selecionada.")
                return
            idx = self._thumb_frames.index(self._selected_thumb)
            self._selected_thumb.destroy()
            self._thumb_frames.pop(idx)
            self._orig_page_idx.pop(idx)
            self._selected_thumb = None
            for i, fr in enumerate(self._thumb_frames):
                fr.grid_forget()
                fr.grid(row=i//3, column=i%3, padx=10, pady=10)
            cvs.configure(scrollregion=cvs.bbox("all"))

        def load_thumbs():
            with fitz.open(self._organize_pdf_path) as doc:
                for orig_idx in range(len(doc)):
                    img = pil_from_page(orig_idx)
                    img.thumbnail((190, 270))
                    tkimg = ImageTk.PhotoImage(img)

                    box = ctk.CTkFrame(inner, fg_color="#2E2E2E", corner_radius=8)
                    box.orig_index = orig_idx

                    il = tk.Label(box, image=tkimg, bg="#2E2E2E")
                    il.pack(padx=6, pady=6)
                    tl = ctk.CTkLabel(box, text=f"P. {orig_idx+1}")
                    tl.pack(pady=(0, 6))

                    def on_click(event, frame=box):
                        if self._selected_thumb:
                            self._selected_thumb.configure(fg_color="#2E2E2E")
                        self._selected_thumb = frame
                        frame.configure(fg_color=ACCENT)
                        show_preview(frame.orig_index)

                    def on_press(event, frame=box):
                        self._dragged_frame = frame
                        self._drag_start_index = self._thumb_frames.index(frame)
                        frame.configure(fg_color=BRAND_BLUE2)

                    def on_drag(event, frame=box):
                        if getattr(self,"_dragged_frame",None) != frame: return
                        x = cvs.winfo_pointerx() - inner.winfo_rootx()
                        y = cvs.winfo_pointery() - inner.winfo_rooty()
                        iw = frame.winfo_width()  + 20
                        ih = frame.winfo_height() + 20
                        col = max(0, min(int(x//iw), 2))
                        row_n = max(0, int(y//ih))
                        tgt   = max(0, min(row_n*3+col, len(self._thumb_frames)-1))
                        if tgt < len(self._thumb_frames):
                            t = self._thumb_frames[tgt]
                            drag_indicator.place(x=t.winfo_x(), y=t.winfo_y()-6,
                                                  width=t.winfo_width())
                        self._drop_target_index = tgt

                    def on_release(event, frame=box):
                        drag_indicator.place_forget()
                        if not hasattr(self, "_drop_target_index"): return
                        src = self._drag_start_index
                        tgt = self._drop_target_index
                        if src != tgt and tgt < len(self._thumb_frames):
                            self._thumb_frames.insert(
                                tgt, self._thumb_frames.pop(src))
                            self._orig_page_idx.insert(
                                tgt, self._orig_page_idx.pop(src))
                        for i, f in enumerate(self._thumb_frames):
                            f.grid_forget()
                            f.grid(row=i//3, column=i%3, padx=10, pady=10)
                            f.configure(fg_color="#2E2E2E")
                        frame.configure(fg_color=ACCENT)
                        self._selected_thumb = frame
                        if hasattr(self, "_drop_target_index"):
                            del self._drop_target_index

                    for w in (box, il, tl):
                        w.bind("<Button-1>",      on_click)
                        w.bind("<ButtonPress-1>",  on_press)
                        w.bind("<B1-Motion>",      on_drag)
                        w.bind("<ButtonRelease-1>",on_release)

                    self._thumb_images.append(tkimg)
                    self._thumb_frames.append(box)
                    self._orig_page_idx.append(orig_idx)
                    g = len(self._thumb_frames) - 1
                    box.grid(row=g//3, column=g%3, padx=10, pady=10)

            cvs.configure(scrollregion=cvs.bbox("all"))

        run_in_thread(load_thumbs)

    def save_organized_pdf(self, win):
        path = filedialog.asksaveasfilename(defaultextension=".pdf",
                                             filetypes=[("PDF", "*.pdf")])
        if not path: return
        try:
            with fitz.open(self._organize_pdf_path) as doc:
                nd = fitz.open()
                for i in self._orig_page_idx:
                    nd.insert_pdf(doc, from_page=i, to_page=i)
                nd.save(path); nd.close()
            messagebox.showinfo("Sucesso", f"PDF salvo:\n{path}")
            win.destroy()
        except Exception as e:
            messagebox.showerror("Erro", str(e))

    def execute_split_visual(self, mode, interval_text, every_text):
        if not self._orig_page_idx:
            messagebox.showerror("Erro", "Nenhuma página carregada."); return
        save_dir = filedialog.askdirectory(title="Pasta de destino")
        if not save_dir: return
        try:
            with fitz.open(self._organize_pdf_path) as doc:
                ordered = self._orig_page_idx[:]
                def save_part(pidxs, part_i):
                    nd = fitz.open()
                    for p in pidxs:
                        nd.insert_pdf(doc, from_page=p, to_page=p)
                    nd.save(os.path.join(save_dir, f"parte_{part_i}.pdf"))
                    nd.close()
                if mode == "single":
                    for i, p in enumerate(ordered, 1):
                        save_part([p], i)
                elif mode == "every":
                    try:
                        step = int(every_text)
                        if step <= 0: raise ValueError
                    except Exception:
                        messagebox.showerror("Erro", "Número inválido."); return
                    for i, start in enumerate(range(0, len(ordered), step), 1):
                        save_part(ordered[start:start+step], i)
                elif mode == "interval":
                    for i, part in enumerate(interval_text.split(","), 1):
                        part = part.strip()
                        if "-" in part:
                            a, b = map(int, part.split("-"))
                            pages = ordered[a-1:b]
                        else:
                            pages = [ordered[int(part)-1]]
                        save_part(pages, i)
            messagebox.showinfo("Concluído", "PDF dividido com sucesso.")
        except Exception as e:
            messagebox.showerror("Erro", f"Falha:\n{e}")

    # ══════════════════════════════════════════════════════════════════════════
    # 🔄 GIRAR PÁGINAS
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_rotate(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "🔄", "Girar Páginas PDF",
                          "Gire páginas individuais ou todas de uma vez.")
        card = self._card(frame)
        card.pack(fill="x", padx=28, pady=14)

        r1 = ctk.CTkFrame(card, fg_color="transparent")
        r1.pack(fill="x", padx=16, pady=(14, 4))
        self.rotate_file_var = ctk.StringVar(value="Nenhum arquivo")
        self._rotate_pdf_path = None
        self._action_btn(r1, "📂 Selecionar PDF",
                          self._rotate_select, width=200).pack(side="left")
        ctk.CTkLabel(r1, textvariable=self.rotate_file_var,
                     text_color=TEXT_DIM).pack(side="left", padx=10)

        ctk.CTkLabel(card,
                     text="Páginas a girar (ex: todas, 1, 2-5, 1,3,5):",
                     anchor="w").pack(anchor="w", padx=16, pady=(8, 2))
        self.rotate_pages_entry = ctk.CTkEntry(card, height=38,
                                                placeholder_text="todas")
        self.rotate_pages_entry.pack(fill="x", padx=16)
        self.rotate_pages_entry.insert(0, "todas")

        af = ctk.CTkFrame(card, fg_color="transparent")
        af.pack(fill="x", padx=16, pady=8)
        ctk.CTkLabel(af, text="Ângulo:").pack(side="left")
        self.rotate_angle_var = ctk.StringVar(value="90")
        for ang in ["90", "180", "270"]:
            ctk.CTkRadioButton(af, text=f"{ang}°",
                                variable=self.rotate_angle_var,
                                value=ang).pack(side="left", padx=12)

        self._action_btn(card, "🔄 Girar e Salvar",
                          self._rotate_run).pack(pady=14)
        self.rotate_prog = ctk.CTkProgressBar(card)
        self.rotate_prog.pack(fill="x", padx=16, pady=(0, 14))
        self.rotate_prog.set(0)

    def _rotate_select(self):
        f = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if f:
            self._rotate_pdf_path = f
            self.rotate_file_var.set(Path(f).name)

    def _rotate_run(self):
        if not self._rotate_pdf_path:
            messagebox.showwarning("Aviso", "Selecione um PDF."); return
        angle     = int(self.rotate_angle_var.get())
        pages_raw = self.rotate_pages_entry.get().strip().lower()
        out = filedialog.asksaveasfilename(defaultextension=".pdf",
                                            filetypes=[("PDF", "*.pdf")],
                                            initialfile="rotacionado.pdf")
        if not out: return

        def worker():
            try:
                reader = PdfReader(self._rotate_pdf_path)
                total  = len(reader.pages)
                if pages_raw in ("todas", "all", "*", ""):
                    target = set(range(total))
                else:
                    target = set()
                    for part in pages_raw.split(","):
                        part = part.strip()
                        if "-" in part:
                            a, b = map(int, part.split("-"))
                            target.update(range(a-1, min(b, total)))
                        else:
                            target.add(int(part)-1)
                writer = PdfWriter()
                for i, page in enumerate(reader.pages):
                    if i in target:
                        page.rotate(angle)
                    writer.add_page(page)
                with open(out, "wb") as fo:
                    writer.write(fo)
                self.rotate_prog.set(1)
                self.set_status(f"Rotacionado: {Path(out).name}", 1)
                self.after(0, lambda: messagebox.showinfo(
                    "Concluído", f"PDF salvo:\n{out}"))
            except Exception as e:
                self.after(0, lambda: messagebox.showerror("Erro", str(e)))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 🖼️ IMAGENS → PDF
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_images(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "🖼️", "Imagens → PDF",
                          "Converta imagens para PDF.")
        card = self._card(frame)
        card.pack(fill="both", expand=True, padx=28, pady=14)

        self.img_listbox = self._file_listbox(card, height=160)
        self.img_listbox.pack(fill="x", padx=16, pady=(14, 4))
        self.img_files: List[str] = []

        self.img_single_pdf = ctk.CTkCheckBox(
            card, text="Gerar 1 PDF único com todas as imagens")
        self.img_single_pdf.pack(anchor="w", padx=16, pady=8)

        self.img_outdir = ctk.StringVar(value=str(Path.home()))
        r2 = ctk.CTkFrame(card, fg_color="transparent")
        r2.pack(fill="x", padx=16)
        self._action_btn(r2, "📁 Pasta de saída",
                          self.select_img_outdir, width=170,
                          color="#1E3A5F").pack(side="left")
        ctk.CTkLabel(r2, textvariable=self.img_outdir,
                     text_color=TEXT_DIM).pack(side="left", padx=8)

        btns = ctk.CTkFrame(card, fg_color="transparent")
        btns.pack(fill="x", padx=16, pady=(8, 4))
        self._action_btn(btns, "➕ Adicionar Imagens",
                          self.select_image_files, width=180
                          ).pack(side="left", padx=(0, 8))
        self._action_btn(btns, "🚀 Converter",
                          self.start_images, width=150).pack(side="right")

        self.img_progress = ctk.CTkProgressBar(card)
        self.img_progress.pack(fill="x", padx=16, pady=(4, 14))
        self.img_progress.set(0)

    def select_image_files(self):
        paths = filedialog.askopenfilenames(
            filetypes=[("Images", "*.png;*.jpg;*.jpeg;*.tiff;*.bmp")])
        if paths:
            self.img_files = list(paths)
            self.img_listbox.delete("1.0", "end")
            for f in self.img_files:
                self.img_listbox.insert("end", f"{Path(f).name}\n")

    def select_img_outdir(self):
        d = filedialog.askdirectory()
        if d: self.img_outdir.set(d)

    def start_images(self):
        files = getattr(self, "img_files", [])
        if not files:
            messagebox.showwarning("Aviso", "Nenhuma imagem selecionada."); return
        outdir = Path(self.img_outdir.get())
        outdir.mkdir(parents=True, exist_ok=True)
        single = bool(self.img_single_pdf.get())

        def worker():
            self.set_status("Convertendo imagens...", 0)
            self.img_progress.set(0)
            try:
                out = images_to_pdfs(files, str(outdir), single_pdf=single)
                self.img_progress.set(1.0)
                self.set_status(f"{len(out)} arquivo(s) gerado(s)", 1)
                self.after(0, lambda: messagebox.showinfo(
                    "Concluído", f"Salvo em:\n{outdir}"))
            except Exception as e:
                self.after(0, lambda: messagebox.showerror("Erro", str(e)))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 📸 PDF → IMAGENS
    # ══════════════════════════════════════════════════════════════════════════

    def open_pdf_to_images(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "📸", "PDF → Imagens",
                          "Converta cada página em PNG ou JPG.")
        card = self._card(frame)
        card.pack(fill="x", padx=28, pady=14)

        r1 = ctk.CTkFrame(card, fg_color="transparent")
        r1.pack(fill="x", padx=16, pady=(14, 4))
        self.p2i_files_var = ctk.StringVar(value="Nenhum arquivo")
        self._p2i_files: List[str] = []
        self._action_btn(r1, "📂 Selecionar PDFs",
                          self._p2i_select, width=200).pack(side="left")
        ctk.CTkLabel(r1, textvariable=self.p2i_files_var,
                     text_color=TEXT_DIM).pack(side="left", padx=10)

        opts = ctk.CTkFrame(card, fg_color="transparent")
        opts.pack(fill="x", padx=16, pady=8)
        ctk.CTkLabel(opts, text="Formato:").pack(side="left")
        self.p2i_fmt = ctk.StringVar(value="png")
        for fmt in ["PNG", "JPG"]:
            ctk.CTkRadioButton(opts, text=fmt, variable=self.p2i_fmt,
                                value=fmt.lower()).pack(side="left", padx=10)

        ctk.CTkLabel(card, text="DPI:", anchor="w").pack(anchor="w", padx=16)
        self.p2i_dpi_var = ctk.IntVar(value=200)
        sl = ctk.CTkSlider(card, from_=72, to=400, number_of_steps=32,
                            variable=self.p2i_dpi_var)
        sl.pack(fill="x", padx=16)
        self.p2i_dpi_lbl = ctk.CTkLabel(card, text="200 DPI", text_color=TEXT_DIM)
        self.p2i_dpi_lbl.pack(anchor="e", padx=16)
        sl.configure(command=lambda v: self.p2i_dpi_lbl.configure(
            text=f"{int(float(v))} DPI"))

        self._action_btn(card, "🚀 Converter", self._p2i_run).pack(pady=14)
        self.p2i_prog = ctk.CTkProgressBar(card)
        self.p2i_prog.pack(fill="x", padx=16, pady=(0, 14))
        self.p2i_prog.set(0)

    def _p2i_select(self):
        fs = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if fs:
            self._p2i_files = list(fs)
            self.p2i_files_var.set(f"{len(fs)} arquivo(s)")

    def _p2i_run(self):
        if not self._p2i_files:
            messagebox.showwarning("Aviso", "Selecione PDFs."); return
        out_dir = filedialog.askdirectory(title="Pasta de saída")
        if not out_dir: return
        fmt, dpi = self.p2i_fmt.get(), self.p2i_dpi_var.get()

        def worker():
            total = len(self._p2i_files)
            for fi, pdf_path in enumerate(self._p2i_files):
                name = Path(pdf_path).stem
                with fitz.open(pdf_path) as doc:
                    for pi, page in enumerate(doc):
                        pix = page.get_pixmap(dpi=dpi)
                        op  = unique_path(
                            Path(out_dir)/f"{name}_P{pi+1}.{fmt}")
                        if fmt == "jpg":
                            pix.save(str(op), output="jpeg")
                        else:
                            pix.save(str(op))
                self.p2i_prog.set((fi+1)/total)
                self.set_status(f"Convertido {fi+1}/{total}", (fi+1)/total)
            self.set_status("Conversão concluída!", 1)
            self.after(0, lambda: messagebox.showinfo(
                "Concluído", f"Imagens salvas em:\n{out_dir}"))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 📝 WORD → PDF
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_word(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "📝", "Word → PDF",
                          "Converta .docx em PDF em lote.")
        card = self._card(frame)
        card.pack(fill="x", padx=28, pady=14)

        self.docx_files: List[str] = []
        self.docx_files_var = ctk.StringVar(value="Nenhum arquivo")
        self.docx_outdir    = ctk.StringVar(value=str(Path.home()))

        r1 = ctk.CTkFrame(card, fg_color="transparent")
        r1.pack(fill="x", padx=16, pady=(14, 4))
        self._action_btn(r1, "📂 Selecionar .docx",
                          self.select_docx_files, width=200).pack(side="left")
        ctk.CTkLabel(r1, textvariable=self.docx_files_var,
                     text_color=TEXT_DIM).pack(side="left", padx=10)

        r2 = ctk.CTkFrame(card, fg_color="transparent")
        r2.pack(fill="x", padx=16, pady=4)
        self._action_btn(r2, "📁 Pasta de saída",
                          self.select_docx_outdir, width=200,
                          color="#1E3A5F").pack(side="left")
        ctk.CTkLabel(r2, textvariable=self.docx_outdir,
                     text_color=TEXT_DIM).pack(side="left", padx=10)

        self._action_btn(card, "🚀 Converter .docx → PDF",
                          self.start_docx).pack(pady=14)
        self.docx_progress = ctk.CTkProgressBar(card)
        self.docx_progress.pack(fill="x", padx=16, pady=(0, 14))
        self.docx_progress.set(0)

    def select_docx_files(self):
        paths = filedialog.askopenfilenames(
            filetypes=[("Word files", "*.docx")])
        if paths:
            self.docx_files = list(paths)
            self.docx_files_var.set(f"{len(paths)} arquivo(s)")

    def select_docx_outdir(self):
        d = filedialog.askdirectory()
        if d: self.docx_outdir.set(d)

    def start_docx(self):
        if not getattr(self, "docx_files", []):
            messagebox.showwarning("Aviso", "Nenhum .docx selecionado."); return
        outdir = Path(self.docx_outdir.get())
        outdir.mkdir(parents=True, exist_ok=True)
        run_in_thread(self._do_docx, self.docx_files, outdir)

    def _do_docx(self, files, outdir: Path):
        self.set_status("Convertendo .docx...", 0)
        self.docx_progress.set(0)
        try:
            for i, file in enumerate(files, 1):
                try:
                    out_file = unique_path(outdir / (Path(file).stem + ".pdf"))
                    if docx2pdf_convert:
                        docx2pdf_convert(file, str(out_file))
                    elif HAVE_DOCX:
                        doc = DocxDocument(file)
                        c   = rl_canvas.Canvas(str(out_file), pagesize=A4)
                        w, h = A4; x, y = 50, h-50
                        for para in doc.paragraphs:
                            text = para.text.strip()
                            if not text: y -= 20; continue
                            c.drawString(x, y, text[:120]); y -= 15
                            if y < 50: c.showPage(); y = h-50
                        c.save()
                    else:
                        soffice = shutil.which("soffice") or shutil.which("libreoffice")
                        if soffice:
                            subprocess.run([soffice, "--headless",
                                            "--convert-to", "pdf",
                                            "--outdir", str(outdir),
                                            str(file)], check=True)
                        else:
                            raise RuntimeError("Nenhum conversor disponível.")
                except Exception as e:
                    print(f"[docx err] {file}: {e}")
                self.docx_progress.set(i / len(files))
            self.docx_progress.set(1.0)
            self.set_status("Conversão .docx concluída.", 1)
            self.after(0, lambda: messagebox.showinfo(
                "Sucesso", f"Arquivos salvos em:\n{outdir}"))
        except Exception as e:
            self.after(0, lambda: messagebox.showerror("Erro", str(e)))

    # ══════════════════════════════════════════════════════════════════════════
    # 📄 PDF → WORD  — com log de erros visível
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_pdf_to_word(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "📄", "PDF → Word",
                          "Converta PDFs em .docx editáveis (OCR quando necessário).")
        card = self._card(frame)
        card.pack(fill="both", expand=True, padx=28, pady=14)

        ctk.CTkLabel(card,
            text="ℹ️  Usa texto nativo + OCR (Tesseract) para PDFs escaneados. "
                 "Erros por página aparecem no log abaixo.",
            wraplength=700, font=ctk.CTkFont("Segoe UI", 12),
            text_color=TEXT_DIM).pack(anchor="w", padx=16, pady=(14, 4))

        self._action_btn(card, "🚀 Selecionar PDFs e Converter",
                          self.run_pdf_to_word_ocr, width=300).pack(pady=14)

        ctk.CTkLabel(card, text="Log de conversão:",
                     anchor="w", font=ctk.CTkFont("Segoe UI", 12)
                     ).pack(anchor="w", padx=16)
        self.p2w_log = ctk.CTkTextbox(
            card, height=220,
            font=ctk.CTkFont("Consolas", 11),
            fg_color="#12121E", text_color="#A0AEC0")
        self.p2w_log.pack(fill="both", expand=True, padx=16, pady=(2, 4))
        self.p2w_log.insert("1.0", "Aguardando conversão...\n")

        self.p2w_prog = ctk.CTkProgressBar(card)
        self.p2w_prog.pack(fill="x", padx=16, pady=(0, 14))
        self.p2w_prog.set(0)

    def run_pdf_to_word_ocr(self):
        if not HAVE_DOCX:
            messagebox.showerror("Erro",
                "Instale python-docx:\npip install python-docx"); return
        pdfs = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if not pdfs: return
        out_dir = filedialog.askdirectory(title="Pasta de saída")
        if not out_dir: return

        def log(msg):
            self.after(0, lambda m=msg: (
                self.p2w_log.insert("end", m + "\n"),
                self.p2w_log.see("end")))

        def worker():
            self.p2w_log.delete("1.0", "end")
            total  = len(pdfs)
            errors = 0
            for i, pdf_path in enumerate(pdfs):
                log(f"[{i+1}/{total}] {Path(pdf_path).name}")
                try:
                    doc  = fitz.open(pdf_path)
                    word = DocxDocument()
                    word.add_heading(Path(pdf_path).stem, 0)
                    for pi, page in enumerate(doc):
                        try:
                            pix = page.get_pixmap(dpi=300)
                            img = Image.open(BytesIO(pix.tobytes("png")))
                            if HAVE_TESSERACT:
                                ocr = pytesseract.image_to_data(
                                    img, lang="por",
                                    output_type=pytesseract.Output.DICT,
                                    config="--psm 6")
                                cur_line = ""; last_y = None
                                for j in range(len(ocr["text"])):
                                    txt = ocr["text"][j].strip()
                                    if not txt: continue
                                    y = ocr["top"][j]
                                    if last_y is not None and abs(y-last_y) > 15:
                                        if cur_line.strip():
                                            word.add_paragraph(cur_line)
                                        cur_line = txt
                                    else:
                                        cur_line += " " + txt
                                    last_y = y
                                if cur_line.strip():
                                    word.add_paragraph(cur_line)
                            else:
                                for ln in page.get_text("text").split("\n"):
                                    if ln.strip():
                                        word.add_paragraph(ln)
                            tmp = Path(tempfile.gettempdir()) / f"_cmat_p{pi}.png"
                            img.save(tmp)
                            word.add_picture(str(tmp), width=DocxInches(6))
                            log(f"   ✔ Página {pi+1}")
                        except Exception as pe:
                            log(f"   ✗ Página {pi+1} — erro: {pe}")
                            errors += 1
                    doc.close()
                    op = unique_path(
                        Path(out_dir) / (Path(pdf_path).stem + ".docx"))
                    word.save(str(op))
                    log(f"   💾 Salvo: {op.name}\n")
                except Exception as e:
                    log(f"   ✗ FALHA TOTAL: {e}\n")
                    errors += 1
                self.p2w_prog.set((i+1)/total)
                self.set_status(f"Convertendo {i+1}/{total}...", (i+1)/total)
            summary = (f"Concluído. {total-errors}/{total} arquivo(s) "
                       f"convertidos." + (f" {errors} com erro(s)." if errors else ""))
            log("\n" + summary)
            self.set_status(summary, 1)
            self.after(0, lambda: messagebox.showinfo("Concluído", summary))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 📊 PDF → EXCEL  — calibrado com tolerância ajustável
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_pdf_to_excel(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "📊", "PDF → Excel",
                          "Extraia tabelas para .xlsx com tolerância ajustável.")
        if not (HAVE_PDFPLUMBER and HAVE_OPENPYXL):
            ctk.CTkLabel(frame,
                text="⚠️  Instale:\npip install pdfplumber openpyxl",
                text_color="#F87171").pack(pady=20)
            return
        card = self._card(frame)
        card.pack(fill="both", expand=True, padx=28, pady=14)

        ctk.CTkLabel(card, text="Modo de extração:",
                     font=ctk.CTkFont("Segoe UI", 13, "bold"),
                     anchor="w").pack(anchor="w", padx=16, pady=(14, 4))
        self.p2e_mode = ctk.StringVar(value="smart")
        for val, lbl in [
            ("tables", "📋 Só tabelas (usa linhas visíveis)"),
            ("text",   "📃 Texto estruturado (linha por célula)"),
            ("smart",  "📊 Inteligente: tabelas + fallback texto"),
        ]:
            ctk.CTkRadioButton(card, text=lbl,
                                variable=self.p2e_mode,
                                value=val).pack(anchor="w", padx=20, pady=2)

        # Tolerância
        tol_row = ctk.CTkFrame(card, fg_color="transparent")
        tol_row.pack(fill="x", padx=16, pady=(8, 4))
        ctk.CTkLabel(tol_row, text="Tolerância vertical (snap_y):").pack(side="left")
        self.p2e_tol_var = ctk.IntVar(value=3)
        ctk.CTkEntry(tol_row, textvariable=self.p2e_tol_var,
                     width=60).pack(side="left", padx=8)
        ctk.CTkLabel(tol_row,
                     text="(aumente para agrupar linhas próximas)",
                     text_color=TEXT_DIM,
                     font=ctk.CTkFont("Segoe UI", 11)).pack(side="left")

        tol_row2 = ctk.CTkFrame(card, fg_color="transparent")
        tol_row2.pack(fill="x", padx=16, pady=(0, 4))
        ctk.CTkLabel(tol_row2, text="Tolerância horizontal (snap_x):").pack(side="left")
        self.p2e_tol_x_var = ctk.IntVar(value=3)
        ctk.CTkEntry(tol_row2, textvariable=self.p2e_tol_x_var,
                     width=60).pack(side="left", padx=8)

        self._action_btn(card, "🚀 Selecionar PDFs e Converter",
                          self._p2e_run, width=300).pack(pady=12)
        self.p2e_log = self._file_listbox(card, height=100)
        self.p2e_log.pack(fill="x", padx=16, pady=(0, 4))
        self.p2e_prog = ctk.CTkProgressBar(card)
        self.p2e_prog.pack(fill="x", padx=16, pady=(0, 14))
        self.p2e_prog.set(0)

    def _p2e_run(self):
        pdfs = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if not pdfs: return
        out_dir = filedialog.askdirectory(title="Pasta de saída")
        if not out_dir: return
        mode  = self.p2e_mode.get()
        tol_y = self.p2e_tol_var.get()
        tol_x = self.p2e_tol_x_var.get()

        def log(msg):
            self.after(0, lambda m=msg: (
                self.p2e_log.insert("end", m + "\n"),
                self.p2e_log.see("end")))

        def clean_row(row):
            return [str(c).strip() if c is not None else "" for c in row]

        def worker():
            self.p2e_log.delete("1.0", "end")
            total = len(pdfs)
            for i, pdf_path in enumerate(pdfs):
                log(f"[{i+1}/{total}] {Path(pdf_path).name}")
                try:
                    wb  = openpyxl.Workbook()
                    n_s = 0
                    tbl_cfg = {
                        "vertical_strategy":   "lines",
                        "horizontal_strategy": "lines",
                        "snap_y_tolerance":    tol_y,
                        "snap_x_tolerance":    tol_x,
                        "join_tolerance":      tol_y,
                        "edge_min_length":     3,
                        "min_words_vertical":  1,
                        "min_words_horizontal":1,
                        "intersection_y_tolerance": tol_y,
                    }
                    with pdfplumber.open(pdf_path) as pdf:
                        for pi, page in enumerate(pdf.pages):
                            extracted = False
                            if mode in ("tables", "smart"):
                                try:
                                    tables = page.extract_tables(tbl_cfg)
                                except Exception:
                                    tables = page.extract_tables()
                                for ti, table in enumerate(tables):
                                    if not table: continue
                                    n_s += 1
                                    ws = wb.create_sheet(
                                        title=f"P{pi+1}_T{ti+1}"[:31])
                                    header = clean_row(table[0]) if table else []
                                    for ci, val in enumerate(header, 1):
                                        cell = ws.cell(row=1, column=ci, value=val)
                                        cell.font = XLFont(bold=True, color="FFFFFF")
                                        cell.fill = PatternFill("solid", fgColor="1B4F8A")
                                    for ri, row_data in enumerate(table[1:], 2):
                                        for ci, val in enumerate(clean_row(row_data), 1):
                                            ws.cell(row=ri, column=ci, value=val)
                                    for col in ws.columns:
                                        mx = max((len(str(c.value)) if c.value else 0)
                                                  for c in col)
                                        ws.column_dimensions[
                                            col[0].column_letter].width = min(mx+4, 50)
                                    extracted = True
                                    log(f"   ✔ P{pi+1} T{ti+1} "
                                        f"({len(table)} linhas)")
                            if not extracted and mode in ("text", "smart"):
                                n_s += 1
                                ws   = wb.create_sheet(title=f"Pag_{pi+1}"[:31])
                                text = page.extract_text() or ""
                                for ln in text.split("\n"):
                                    ws.append([ln])
                                log(f"   ✔ P{pi+1} — texto")
                    if "Sheet" in wb.sheetnames and n_s > 0:
                        del wb["Sheet"]
                    op = unique_path(
                        Path(out_dir) / (Path(pdf_path).stem + ".xlsx"))
                    wb.save(str(op))
                    log(f"   💾 {op.name}")
                except Exception as e:
                    log(f"   ✗ Erro: {e}")
                self.p2e_prog.set((i+1)/total)
            self.set_status("Excel gerado!", 1)
            self.after(0, lambda: messagebox.showinfo(
                "Concluído", f"Planilhas em:\n{out_dir}"))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 🖥️ PDF → POWERPOINT
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_pdf_to_pptx(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "🖥️", "PDF → PowerPoint",
                          "Cada página vira um slide editável.")
        if not HAVE_PPTX:
            ctk.CTkLabel(frame,
                text="⚠️  Instale: pip install python-pptx",
                text_color="#F87171").pack(pady=20)
            return
        card = self._card(frame)
        card.pack(fill="x", padx=28, pady=14)

        ctk.CTkLabel(card, text="DPI da renderização:", anchor="w"
                     ).pack(anchor="w", padx=16, pady=(14, 2))
        self.p2p_dpi_var = ctk.IntVar(value=150)
        sl = ctk.CTkSlider(card, from_=72, to=300, number_of_steps=23,
                            variable=self.p2p_dpi_var)
        sl.pack(fill="x", padx=16)
        self.p2p_dpi_lbl = ctk.CTkLabel(card, text="150 DPI",
                                          text_color=TEXT_DIM)
        self.p2p_dpi_lbl.pack(anchor="e", padx=16)
        sl.configure(command=lambda v: self.p2p_dpi_lbl.configure(
            text=f"{int(float(v))} DPI"))

        self._action_btn(card, "🚀 Selecionar PDFs e Converter",
                          self._p2p_run, width=300).pack(pady=14)
        self.p2p_prog = ctk.CTkProgressBar(card)
        self.p2p_prog.pack(fill="x", padx=16, pady=(0, 14))
        self.p2p_prog.set(0)

    def _p2p_run(self):
        pdfs = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if not pdfs: return
        out_dir = filedialog.askdirectory(title="Pasta de saída")
        if not out_dir: return
        dpi = self.p2p_dpi_var.get()

        def worker():
            total = len(pdfs)
            for i, pdf_path in enumerate(pdfs):
                try:
                    prs = Presentation()
                    prs.slide_width  = Inches(13.33)
                    prs.slide_height = Inches(7.5)
                    blank = prs.slide_layouts[6]
                    with fitz.open(pdf_path) as doc:
                        n = len(doc)
                        for pi, page in enumerate(doc):
                            pix   = page.get_pixmap(dpi=dpi)
                            img   = Image.open(BytesIO(pix.tobytes("png")))
                            buf   = BytesIO()
                            img.save(buf, "PNG"); buf.seek(0)
                            slide = prs.slides.add_slide(blank)
                            slide.shapes.add_picture(
                                buf, Inches(0), Inches(0),
                                prs.slide_width, prs.slide_height)
                            self.set_status(
                                f"PDF {i+1}/{total} — slide {pi+1}/{n}",
                                (i + pi/n) / total)
                    op = unique_path(
                        Path(out_dir) / (Path(pdf_path).stem + ".pptx"))
                    prs.save(str(op))
                except Exception as e:
                    print("p2p err:", e)
                self.p2p_prog.set((i+1)/total)
            self.set_status("PowerPoint gerado!", 1)
            self.after(0, lambda: messagebox.showinfo(
                "Concluído", f"Apresentações em:\n{out_dir}"))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 🔒 PROTEGER COM SENHA
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_protect(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "🔒", "Proteger com Senha",
                          "Criptografe PDFs com senha de abertura e de proprietário.")
        card = self._card(frame)
        card.pack(fill="x", padx=28, pady=14)

        self.prot_files: List[str] = []
        self.prot_files_var = ctk.StringVar(value="Nenhum arquivo")
        r1 = ctk.CTkFrame(card, fg_color="transparent")
        r1.pack(fill="x", padx=16, pady=(14, 4))
        self._action_btn(r1, "📂 Selecionar PDFs",
                          self._prot_select, width=200).pack(side="left")
        ctk.CTkLabel(r1, textvariable=self.prot_files_var,
                     text_color=TEXT_DIM).pack(side="left", padx=10)

        ctk.CTkLabel(card, text="Senha do usuário (para abrir):",
                     anchor="w").pack(anchor="w", padx=16, pady=(12, 2))
        self.prot_user_pw = ctk.CTkEntry(card, show="●", height=38,
                                          placeholder_text="Senha de abertura")
        self.prot_user_pw.pack(fill="x", padx=16)

        ctk.CTkLabel(card,
                     text="Senha do proprietário (para editar/imprimir — opcional):",
                     anchor="w").pack(anchor="w", padx=16, pady=(8, 2))
        self.prot_owner_pw = ctk.CTkEntry(card, show="●", height=38,
                                           placeholder_text="Deixe em branco para usar a mesma")
        self.prot_owner_pw.pack(fill="x", padx=16)

        ctk.CTkLabel(card, text="Permissões:",
                     font=ctk.CTkFont("Segoe UI", 13, "bold"),
                     anchor="w").pack(anchor="w", padx=16, pady=(10, 2))
        pf = ctk.CTkFrame(card, fg_color="transparent")
        pf.pack(fill="x", padx=16, pady=(0, 14))
        self.perm_print  = ctk.BooleanVar(value=True)
        self.perm_copy   = ctk.BooleanVar(value=True)
        self.perm_modify = ctk.BooleanVar(value=False)
        for txt, var in [("Permitir impressão", self.perm_print),
                          ("Permitir cópia de texto", self.perm_copy),
                          ("Permitir modificação", self.perm_modify)]:
            ctk.CTkCheckBox(pf, text=txt, variable=var).pack(anchor="w", pady=2)

        self._action_btn(card, "🔒 Proteger PDFs",
                          self._prot_run).pack(pady=(0, 14))
        self.prot_prog = ctk.CTkProgressBar(card)
        self.prot_prog.pack(fill="x", padx=16, pady=(0, 14))
        self.prot_prog.set(0)

    def _prot_select(self):
        fs = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if fs:
            self.prot_files = list(fs)
            self.prot_files_var.set(f"{len(fs)} arquivo(s)")

    def _prot_run(self):
        if not self.prot_files:
            messagebox.showwarning("Aviso", "Selecione PDFs."); return
        user_pw  = self.prot_user_pw.get().strip()
        owner_pw = self.prot_owner_pw.get().strip() or user_pw
        if not user_pw:
            messagebox.showwarning("Aviso", "Digite a senha do usuário."); return
        out_dir = filedialog.askdirectory(title="Pasta de saída")
        if not out_dir: return

        def worker():
            total = len(self.prot_files)
            for i, fp in enumerate(self.prot_files):
                try:
                    reader = PdfReader(fp); writer = PdfWriter()
                    for page in reader.pages:
                        writer.add_page(page)
                    perms = 0
                    if self.perm_print.get():  perms |= 4
                    if self.perm_copy.get():   perms |= 16
                    if self.perm_modify.get(): perms |= 8
                    writer.encrypt(user_pw, owner_pw, permissions_flag=perms)
                    op = unique_path(
                        Path(out_dir) / (Path(fp).stem + "_protegido.pdf"))
                    with open(op, "wb") as fo:
                        writer.write(fo)
                except Exception as e:
                    print("prot err:", e)
                self.prot_prog.set((i+1)/total)
            self.set_status("Proteção aplicada!", 1)
            self.after(0, lambda: messagebox.showinfo(
                "Concluído", f"PDFs protegidos em:\n{out_dir}"))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 🔓 REMOVER SENHA
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_remove_password(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "🔓", "Remover Senha de PDF",
                          "Descriptografe PDFs protegidos.")
        card = self._card(frame)
        card.pack(fill="x", padx=28, pady=14)

        self.unprot_files: List[str] = []
        self.unprot_files_var = ctk.StringVar(value="Nenhum arquivo")
        r1 = ctk.CTkFrame(card, fg_color="transparent")
        r1.pack(fill="x", padx=16, pady=(14, 4))
        self._action_btn(r1, "📂 Selecionar PDFs",
                          self._unprot_select, width=200).pack(side="left")
        ctk.CTkLabel(r1, textvariable=self.unprot_files_var,
                     text_color=TEXT_DIM).pack(side="left", padx=10)

        ctk.CTkLabel(card, text="Senha atual:", anchor="w"
                     ).pack(anchor="w", padx=16, pady=(10, 2))
        self.unprot_pw = ctk.CTkEntry(card, show="●", height=38,
                                       placeholder_text="Digite a senha")
        self.unprot_pw.pack(fill="x", padx=16)

        self._action_btn(card, "🔓 Remover Senha",
                          self._unprot_run).pack(pady=14)
        self.unprot_prog = ctk.CTkProgressBar(card)
        self.unprot_prog.pack(fill="x", padx=16, pady=(0, 14))
        self.unprot_prog.set(0)

    def _unprot_select(self):
        fs = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if fs:
            self.unprot_files = list(fs)
            self.unprot_files_var.set(f"{len(fs)} arquivo(s)")

    def _unprot_run(self):
        if not self.unprot_files:
            messagebox.showwarning("Aviso", "Selecione PDFs."); return
        pw      = self.unprot_pw.get()
        out_dir = filedialog.askdirectory(title="Pasta de saída")
        if not out_dir: return

        def worker():
            total = len(self.unprot_files); errors = 0
            for i, fp in enumerate(self.unprot_files):
                try:
                    reader = PdfReader(fp)
                    if reader.is_encrypted:
                        ok = reader.decrypt(pw)
                        if not ok:
                            errors += 1
                            print(f"Senha incorreta: {fp}"); continue
                    writer = PdfWriter()
                    for page in reader.pages:
                        writer.add_page(page)
                    op = unique_path(
                        Path(out_dir) / (Path(fp).stem + "_sem_senha.pdf"))
                    with open(op, "wb") as fo:
                        writer.write(fo)
                except Exception as e:
                    print("unprot err:", e); errors += 1
                self.unprot_prog.set((i+1)/total)
            msg = (f"Senha removida de {total-errors}/{total} arquivo(s).")
            self.set_status(msg, 1)
            self.after(0, lambda: messagebox.showinfo(
                "Concluído", msg + f"\nDestino: {out_dir}"))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 💧 MARCA D'ÁGUA  — com posicionamento interativo por clique no preview
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_watermark(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "💧", "Marca D'água",
                          "Aplique texto ou imagem. "
                          "Clique no preview para posicionar.")

        card = self._card(frame)
        card.pack(fill="both", expand=True, padx=28, pady=14)

        cols = ctk.CTkFrame(card, fg_color="transparent")
        cols.pack(fill="both", expand=True, padx=8, pady=8)
        cols.grid_columnconfigure(0, weight=1)
        cols.grid_columnconfigure(1, weight=0)

        ctrl = ctk.CTkFrame(cols, fg_color="transparent")
        ctrl.grid(row=0, column=0, sticky="nsew", padx=(8, 4))

        # Arquivos alvo
        self.wm_files: List[str] = []
        self.wm_files_var = ctk.StringVar(value="Nenhum PDF")
        r1 = ctk.CTkFrame(ctrl, fg_color="transparent")
        r1.pack(fill="x", pady=(4, 2))
        self._action_btn(r1, "📂 PDFs alvo",
                          self._wm_select_files, width=150).pack(side="left")
        ctk.CTkLabel(r1, textvariable=self.wm_files_var,
                     text_color=TEXT_DIM).pack(side="left", padx=8)

        # Imagem
        self.wm_img_path = None
        self.wm_img_var  = ctk.StringVar(value="Nenhuma imagem")
        r2 = ctk.CTkFrame(ctrl, fg_color="transparent")
        r2.pack(fill="x", pady=2)
        self._action_btn(r2, "🖼️ Imagem da marca",
                          self._wm_select_img, width=170,
                          color="#1E3A5F").pack(side="left")
        ctk.CTkLabel(r2, textvariable=self.wm_img_var,
                     text_color=TEXT_DIM).pack(side="left", padx=8)

        # Texto
        ctk.CTkLabel(ctrl, text="Texto da marca (opcional):",
                     anchor="w").pack(anchor="w", pady=(6, 2))
        self.wm_text_entry = ctk.CTkEntry(ctrl, height=36,
                                           placeholder_text="ex: CONFIDENCIAL")
        self.wm_text_entry.pack(fill="x")

        # Escala
        sz_r = ctk.CTkFrame(ctrl, fg_color="transparent")
        sz_r.pack(fill="x", pady=(8, 0))
        ctk.CTkLabel(sz_r, text="Tamanho:").pack(side="left")
        self.wm_scale_var = ctk.DoubleVar(value=1.0)
        self.wm_scale_lbl = ctk.CTkLabel(sz_r, text="100%", text_color=TEXT_DIM)
        self.wm_scale_lbl.pack(side="right")
        sl1 = ctk.CTkSlider(ctrl, from_=0.1, to=3.0, number_of_steps=29,
                              variable=self.wm_scale_var)
        sl1.pack(fill="x")
        sl1.configure(command=lambda v: self.wm_scale_lbl.configure(
            text=f"{int(float(v)*100)}%"))

        # Transparência
        al_r = ctk.CTkFrame(ctrl, fg_color="transparent")
        al_r.pack(fill="x", pady=(4, 0))
        ctk.CTkLabel(al_r, text="Transparência:").pack(side="left")
        self.wm_alpha_var = ctk.DoubleVar(value=0.4)
        self.wm_alpha_lbl = ctk.CTkLabel(al_r, text="40%", text_color=TEXT_DIM)
        self.wm_alpha_lbl.pack(side="right")
        sl2 = ctk.CTkSlider(ctrl, from_=0.05, to=1.0, number_of_steps=19,
                              variable=self.wm_alpha_var)
        sl2.pack(fill="x")
        sl2.configure(command=lambda v: self.wm_alpha_lbl.configure(
            text=f"{int(float(v)*100)}%"))

        # Posição
        self.wm_pos_x = ctk.DoubleVar(value=0.5)
        self.wm_pos_y = ctk.DoubleVar(value=0.5)
        self._wm_pos_lbl = ctk.CTkLabel(
            ctrl, text="📍 Clique no preview para definir a posição",
            text_color=TEXT_DIM, anchor="w")
        self._wm_pos_lbl.pack(anchor="w", pady=(8, 0))

        self._action_btn(ctrl, "🔍 Carregar Preview",
                          self._wm_load_preview, width=220,
                          color="#1E3A5F").pack(pady=(8, 4))
        self._action_btn(ctrl, "💧 Aplicar Marca D'água",
                          self.start_watermark, width=220).pack(pady=(4, 0))
        self.wm_progress = ctk.CTkProgressBar(ctrl)
        self.wm_progress.pack(fill="x", pady=(8, 0))
        self.wm_progress.set(0)

        # Preview interativo
        pv_frame = ctk.CTkFrame(cols, fg_color="#12121E",
                                 width=280, corner_radius=10)
        pv_frame.grid(row=0, column=1, sticky="nsew", padx=(4, 8))
        pv_frame.grid_propagate(False)
        ctk.CTkLabel(pv_frame, text="Preview — clique para posicionar",
                     font=ctk.CTkFont("Segoe UI", 11),
                     text_color=TEXT_DIM).pack(pady=(8, 2))
        self.wm_preview_cv = tk.Canvas(pv_frame, bg="#12121E",
                                        cursor="crosshair", highlightthickness=0)
        self.wm_preview_cv.pack(fill="both", expand=True, padx=4, pady=4)
        self._wm_pdf_preview_path = None

        def on_wm_click(ev):
            cw = self.wm_preview_cv.winfo_width()
            ch = self.wm_preview_cv.winfo_height()
            if cw < 2 or ch < 2: return
            px = ev.x / cw
            py = 1.0 - ev.y / ch     # Y invertido (PDF coords)
            self.wm_pos_x.set(round(px, 3))
            self.wm_pos_y.set(round(py, 3))
            self._wm_pos_lbl.configure(
                text=f"📍 Posição: {int(px*100)}% horiz, "
                     f"{int((1-py)*100)}% vert (do topo)")
            self.wm_preview_cv.delete("pos_marker")
            r = 9
            self.wm_preview_cv.create_oval(
                ev.x-r, ev.y-r, ev.x+r, ev.y+r,
                outline=ACCENT, width=2, tags="pos_marker")
            self.wm_preview_cv.create_line(
                ev.x-16, ev.y, ev.x+16, ev.y,
                fill=ACCENT, width=1, tags="pos_marker")
            self.wm_preview_cv.create_line(
                ev.x, ev.y-16, ev.x, ev.y+16,
                fill=ACCENT, width=1, tags="pos_marker")

        self.wm_preview_cv.bind("<Button-1>", on_wm_click)

    def _wm_select_files(self):
        fs = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if fs:
            self.wm_files = list(fs)
            self.wm_files_var.set(f"{len(fs)} PDF(s)")
            self._wm_pdf_preview_path = self.wm_files[0]
            self._wm_load_preview()

    def _wm_select_img(self):
        f = filedialog.askopenfilename(
            filetypes=[("Imagens", "*.png *.jpg *.jpeg")])
        if f:
            self.wm_img_path = f
            self.wm_img_var.set(Path(f).name)

    def _wm_load_preview(self):
        pdf = self._wm_pdf_preview_path or (
            self.wm_files[0] if self.wm_files else None)
        if not pdf:
            messagebox.showwarning("Aviso", "Selecione PDFs alvo primeiro.")
            return
        self._wm_pdf_preview_path = pdf

        def _r():
            with fitz.open(pdf) as doc:
                pix = doc[0].get_pixmap(dpi=100)
            img = Image.open(BytesIO(pix.tobytes("png")))
            cw  = self.wm_preview_cv.winfo_width()  or 260
            ch  = self.wm_preview_cv.winfo_height() or 370
            r   = min(cw/img.width, ch/img.height)
            img = img.resize((max(1, int(img.width*r)),
                              max(1, int(img.height*r))), Image.LANCZOS)
            tkimg = ImageTk.PhotoImage(img)
            def _u():
                self.wm_preview_cv.delete("all")
                self.wm_preview_cv.create_image(
                    cw//2, ch//2, image=tkimg, anchor="center")
                self.wm_preview_cv.image = tkimg
            self.wm_preview_cv.after(0, _u)
        run_in_thread(_r)

    def start_watermark(self):
        if not self.wm_files:
            messagebox.showwarning("Aviso", "Selecione PDFs alvo."); return
        text  = self.wm_text_entry.get().strip()
        img   = self.wm_img_path
        if not text and not img:
            messagebox.showwarning("Aviso", "Defina texto ou imagem."); return
        scale = self.wm_scale_var.get()
        alpha = self.wm_alpha_var.get()
        pos_x = self.wm_pos_x.get()
        pos_y = self.wm_pos_y.get()
        out_dir = filedialog.askdirectory(title="Pasta de saída")
        if not out_dir: return

        def worker():
            n = len(self.wm_files)
            self.wm_progress.set(0)
            for i, f in enumerate(self.wm_files, 1):
                try:
                    op = unique_path(Path(out_dir) / f"WM_{Path(f).name}")
                    apply_watermark(
                        pdf_path=f,
                        watermark_image=img,
                        watermark_text=text,
                        output_path=str(op),
                        scale_factor=scale,
                        alpha=alpha,
                        pos_x_pct=pos_x,
                        pos_y_pct=pos_y)
                except Exception as e:
                    print(f"wm err {f}: {e}")
                self.wm_progress.set(i/n)
                self.set_status(f"Aplicando {i}/{n}...", i/n)
            self.set_status("Marca d'água aplicada!", 1)
            self.after(0, lambda: messagebox.showinfo(
                "Concluído", f"PDFs salvos em:\n{out_dir}"))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # ✏️ EDITAR PDF  — posicionamento interativo por clique no preview
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_edit_pdf(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "✏️", "Editar PDF",
                          "Adicione texto, imagens ou anotações. "
                          "Clique no preview para posicionar.")

        card = self._card(frame)
        card.pack(fill="both", expand=True, padx=28, pady=14)

        cols = ctk.CTkFrame(card, fg_color="transparent")
        cols.pack(fill="both", expand=True, padx=8, pady=8)
        cols.grid_columnconfigure(0, weight=1)
        cols.grid_columnconfigure(1, weight=0)

        ctrl = ctk.CTkFrame(cols, fg_color="transparent")
        ctrl.grid(row=0, column=0, sticky="nsew", padx=(8, 4))

        # Arquivo
        r1 = ctk.CTkFrame(ctrl, fg_color="transparent")
        r1.pack(fill="x", pady=(4, 2))
        self.edit_file_var = ctk.StringVar(value="Nenhum arquivo")
        self._edit_pdf_path = None
        self._action_btn(r1, "📂 Selecionar PDF",
                          self._edit_select, width=190).pack(side="left")
        ctk.CTkLabel(r1, textvariable=self.edit_file_var,
                     text_color=TEXT_DIM).pack(side="left", padx=8)

        # Página + botão preview
        pg_r = ctk.CTkFrame(ctrl, fg_color="transparent")
        pg_r.pack(fill="x", pady=(4, 2))
        ctk.CTkLabel(pg_r, text="Página:").pack(side="left")
        self.edit_page_entry = ctk.CTkEntry(pg_r, width=60,
                                              placeholder_text="1")
        self.edit_page_entry.pack(side="left", padx=8)
        self._action_btn(pg_r, "🔍 Carregar Preview",
                          self._edit_load_preview, width=180,
                          color="#1E3A5F").pack(side="left")

        # Tipo de edição
        ctk.CTkLabel(ctrl, text="Tipo de edição:",
                     font=ctk.CTkFont("Segoe UI", 12, "bold"),
                     anchor="w").pack(anchor="w", pady=(8, 2))
        self.edit_mode = ctk.StringVar(value="text")
        mf = ctk.CTkFrame(ctrl, fg_color="transparent")
        mf.pack(fill="x")
        for val, lbl in [("text", "📝 Texto"),
                          ("image", "🖼️ Imagem"),
                          ("annotation", "💬 Anotação")]:
            ctk.CTkRadioButton(mf, text=lbl, variable=self.edit_mode,
                                value=val,
                                command=self._edit_mode_changed).pack(
                side="left", padx=6)

        self.edit_opts_frame = ctk.CTkFrame(ctrl, fg_color="transparent")
        self.edit_opts_frame.pack(fill="x", pady=(4, 0))
        self._edit_mode_changed()

        # Posição (definida pelo clique no canvas)
        self._edit_click_x = ctk.DoubleVar(value=100)
        self._edit_click_y = ctk.DoubleVar(value=100)
        self._edit_pos_lbl = ctk.CTkLabel(
            ctrl, text="📍 Clique no preview para definir a posição",
            text_color=TEXT_DIM, anchor="w")
        self._edit_pos_lbl.pack(anchor="w", pady=(6, 0))

        self._action_btn(ctrl, "✏️ Aplicar e Salvar",
                          self._edit_run, width=220).pack(pady=14)
        self.edit_prog = ctk.CTkProgressBar(ctrl)
        self.edit_prog.pack(fill="x")
        self.edit_prog.set(0)

        # ── Preview interativo (coluna direita) ──
        pv_frame = ctk.CTkFrame(cols, fg_color="#12121E",
                                 width=295, corner_radius=10)
        pv_frame.grid(row=0, column=1, sticky="nsew", padx=(4, 8))
        pv_frame.grid_propagate(False)
        ctk.CTkLabel(pv_frame, text="Preview — clique para posicionar",
                     font=ctk.CTkFont("Segoe UI", 11),
                     text_color=TEXT_DIM).pack(pady=(8, 2))
        self.edit_preview_cv = tk.Canvas(
            pv_frame, bg="#12121E",
            cursor="crosshair", highlightthickness=0)
        self.edit_preview_cv.pack(fill="both", expand=True, padx=4, pady=4)

        # atributos de transformação canvas↔PDF
        self._edit_preview_ratio = 1.0
        self._edit_preview_offx  = 0
        self._edit_preview_offy  = 0

        def on_edit_click(ev):
            if not self._edit_pdf_path: return
            try:
                with fitz.open(self._edit_pdf_path) as d:
                    pg   = int(self.edit_page_entry.get() or "1") - 1
                    rect = d[pg].rect
                    pdf_w, pdf_h = rect.width, rect.height
            except Exception:
                pdf_w, pdf_h = 595, 842
            ratio = self._edit_preview_ratio
            ox    = self._edit_preview_offx
            oy    = self._edit_preview_offy
            px = (ev.x - ox) / ratio
            py = pdf_h - (ev.y - oy) / ratio   # Y invertido
            self._edit_click_x.set(round(px, 1))
            self._edit_click_y.set(round(py, 1))
            self._edit_pos_lbl.configure(
                text=f"📍 Posição: x={px:.0f} pt, y={py:.0f} pt")
            self.edit_preview_cv.delete("pos_marker")
            r = 8
            self.edit_preview_cv.create_oval(
                ev.x-r, ev.y-r, ev.x+r, ev.y+r,
                outline=ACCENT, width=2, tags="pos_marker")
            self.edit_preview_cv.create_line(
                ev.x-14, ev.y, ev.x+14, ev.y,
                fill=ACCENT, width=1, tags="pos_marker")
            self.edit_preview_cv.create_line(
                ev.x, ev.y-14, ev.x, ev.y+14,
                fill=ACCENT, width=1, tags="pos_marker")

        self.edit_preview_cv.bind("<Button-1>", on_edit_click)

    def _edit_select(self):
        f = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if f:
            self._edit_pdf_path = f
            self.edit_file_var.set(Path(f).name)
            self._edit_load_preview()

    def _edit_load_preview(self):
        if not self._edit_pdf_path: return
        try:
            pg = int(self.edit_page_entry.get() or "1") - 1
        except ValueError:
            pg = 0

        def _r():
            with fitz.open(self._edit_pdf_path) as doc:
                pix   = doc[pg].get_pixmap(dpi=120)
                pdf_w = doc[pg].rect.width
                pdf_h = doc[pg].rect.height
            img = Image.open(BytesIO(pix.tobytes("png")))
            cw  = self.edit_preview_cv.winfo_width()  or 280
            ch  = self.edit_preview_cv.winfo_height() or 400
            r   = min(cw/img.width, ch/img.height)
            nw, nh = max(1, int(img.width*r)), max(1, int(img.height*r))
            img   = img.resize((nw, nh), Image.LANCZOS)
            tkimg = ImageTk.PhotoImage(img)
            ox = (cw - nw) // 2
            oy = (ch - nh) // 2
            self._edit_preview_ratio = nw / pdf_w
            self._edit_preview_offx  = ox
            self._edit_preview_offy  = oy
            def _u():
                self.edit_preview_cv.delete("all")
                self.edit_preview_cv.create_image(
                    cw//2, ch//2, image=tkimg, anchor="center")
                self.edit_preview_cv.image = tkimg
            self.edit_preview_cv.after(0, _u)
        run_in_thread(_r)

    def _edit_mode_changed(self):
        for w in self.edit_opts_frame.winfo_children():
            w.destroy()
        mode = self.edit_mode.get()
        if mode == "text":
            ctk.CTkLabel(self.edit_opts_frame, text="Texto:",
                         anchor="w").pack(anchor="w")
            self.edit_text_entry = ctk.CTkEntry(
                self.edit_opts_frame, height=36,
                placeholder_text="Digite o texto...")
            self.edit_text_entry.pack(fill="x", pady=(2, 4))
            row = ctk.CTkFrame(self.edit_opts_frame, fg_color="transparent")
            row.pack(fill="x")
            ctk.CTkLabel(row, text="Tamanho:").pack(side="left")
            self.edit_fontsize_var = ctk.IntVar(value=14)
            ctk.CTkEntry(row, textvariable=self.edit_fontsize_var,
                         width=55).pack(side="left", padx=4)
            ctk.CTkLabel(row, text="Cor:").pack(side="left", padx=(8, 2))
            self.edit_color_var = ctk.StringVar(value="#000000")
            self._edit_color_btn = ctk.CTkButton(
                row, text="  ", width=40, height=28,
                fg_color=self.edit_color_var.get(),
                command=self._pick_edit_color)
            self._edit_color_btn.pack(side="left")
        elif mode == "image":
            r = ctk.CTkFrame(self.edit_opts_frame, fg_color="transparent")
            r.pack(fill="x")
            self.edit_ins_img_var = ctk.StringVar(value="Nenhuma imagem")
            self._insert_img_path = None
            self._action_btn(r, "🖼️ Escolher imagem",
                              self._pick_insert_img, width=180,
                              color="#1E3A5F").pack(side="left")
            ctk.CTkLabel(r, textvariable=self.edit_ins_img_var,
                         text_color=TEXT_DIM).pack(side="left", padx=8)
            sz_r = ctk.CTkFrame(self.edit_opts_frame, fg_color="transparent")
            sz_r.pack(fill="x", pady=4)
            ctk.CTkLabel(sz_r, text="Largura (pt):").pack(side="left")
            self.edit_img_w = ctk.CTkEntry(sz_r, width=70,
                                            placeholder_text="200")
            self.edit_img_w.pack(side="left", padx=6)
            ctk.CTkLabel(sz_r, text="Altura (pt):").pack(side="left")
            self.edit_img_h = ctk.CTkEntry(sz_r, width=70,
                                            placeholder_text="200")
            self.edit_img_h.pack(side="left", padx=4)
        else:  # annotation
            ctk.CTkLabel(self.edit_opts_frame, text="Texto da anotação:",
                         anchor="w").pack(anchor="w")
            self.edit_annot_entry = ctk.CTkTextbox(self.edit_opts_frame,
                                                    height=70)
            self.edit_annot_entry.pack(fill="x", pady=(2, 4))

    def _pick_edit_color(self):
        res = colorchooser.askcolor(color=self.edit_color_var.get())
        if res and res[1]:
            self.edit_color_var.set(res[1])
            self._edit_color_btn.configure(fg_color=res[1])

    def _pick_insert_img(self):
        f = filedialog.askopenfilename(
            filetypes=[("Imagens", "*.png *.jpg *.jpeg")])
        if f:
            self._insert_img_path = f
            self.edit_ins_img_var.set(Path(f).name)

    def _edit_run(self):
        if not self._edit_pdf_path:
            messagebox.showwarning("Aviso", "Selecione um PDF."); return
        mode = self.edit_mode.get()
        try:
            page_num = int(self.edit_page_entry.get() or "1") - 1
            x = self._edit_click_x.get()
            y = self._edit_click_y.get()
        except ValueError:
            messagebox.showwarning("Aviso", "Página inválida."); return
        out = filedialog.asksaveasfilename(
            defaultextension=".pdf", filetypes=[("PDF", "*.pdf")],
            initialfile=Path(self._edit_pdf_path).stem + "_editado.pdf")
        if not out: return

        def worker():
            self.set_status("Editando PDF...", 0.3)
            try:
                doc  = fitz.open(self._edit_pdf_path)
                page = doc[page_num]
                if mode == "text":
                    txt   = self.edit_text_entry.get()
                    fsize = self.edit_fontsize_var.get()
                    clr   = self.edit_color_var.get().lstrip("#")
                    r, g, b = (int(clr[i:i+2], 16)/255 for i in (0, 2, 4))
                    page.insert_text((x, y), txt, fontsize=fsize,
                                     color=(r, g, b))
                elif mode == "image":
                    if not getattr(self, "_insert_img_path", None):
                        messagebox.showwarning("Aviso",
                            "Escolha uma imagem."); return
                    iw = float(self.edit_img_w.get() or 200)
                    ih = float(self.edit_img_h.get() or 200)
                    rect = fitz.Rect(x, y-ih, x+iw, y)
                    page.insert_image(rect,
                                       filename=self._insert_img_path)
                else:
                    txt = self.edit_annot_entry.get("1.0", "end").strip()
                    page.add_text_annot((x, y), txt)
                doc.save(out); doc.close()
                self.edit_prog.set(1)
                self.set_status(f"Editado: {Path(out).name}", 1)
                self.after(0, lambda: messagebox.showinfo(
                    "Concluído", f"PDF salvo em:\n{out}"))
            except Exception as e:
                self.after(0, lambda: messagebox.showerror("Erro", str(e)))
        run_in_thread(worker)

# ══════════════════════════════════════════════════════════════════════════
    # ✍️ ASSINATURA DIGITAL — preview vertical + múltiplos pontos
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_signature(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(
            frame,
            "✍️",
            "Assinatura Digital",
            "Selecione PDFs, carregue o preview vertical, clique para marcar os pontos e aplique a imagem várias vezes."
        )

        card = self._card(frame)
        card.pack(fill="both", expand=True, padx=28, pady=14)

        cols = ctk.CTkFrame(card, fg_color="transparent")
        cols.pack(fill="both", expand=True, padx=8, pady=8)
        
        cols.grid_columnconfigure(0, weight=1)
        cols.grid_columnconfigure(1, weight=2) 

        ctrl = ctk.CTkFrame(cols, fg_color="transparent")
        ctrl.grid(row=0, column=0, sticky="nsew", padx=(8, 4))

        # Variáveis de Estado
        self.sig_file_var = ctk.StringVar(value="Nenhum PDF selecionado")
        self.sig_img_var = ctk.StringVar(value="Nenhuma imagem")
        self._sig_pdf_paths: List[str] = []
        self._sig_pdf_preview_path = None
        self._sig_positions: List[Tuple[float, float]] = []
        self._sig_pil_base = None      # Guarda a imagem PIL da assinatura carregada
        self._preview_tkimgs = []     # Evita o garbage collector de apagar as assinaturas do canvas

        # Linha 1: Seleção de Arquivos
        r1 = ctk.CTkFrame(ctrl, fg_color="transparent")
        r1.pack(fill="x", pady=(4, 2))

        self._action_btn(
            r1, "📂 Selecionar PDFs", self._sig_select_pdf, width=190
        ).pack(side="left")

        ctk.CTkLabel(
            r1, textvariable=self.sig_file_var, text_color=TEXT_DIM
        ).pack(side="left", padx=8)

        # Linha 2: Importar imagem da assinatura
        r2 = ctk.CTkFrame(ctrl, fg_color="transparent")
        r2.pack(fill="x", pady=(6, 2))

        self._action_btn(
            r2, "🖼️ Importar imagem da assinatura", self._sig_import_img,
            width=240, color="#1E3A5F"
        ).pack(side="left")

        ctk.CTkLabel(
            r2, textvariable=self.sig_img_var, text_color=TEXT_DIM
        ).pack(side="left", padx=8)

        # NOVO: Linha de ajuste de tamanho da assinatura
        r_size = ctk.CTkFrame(ctrl, fg_color="transparent")
        r_size.pack(fill="x", pady=(8, 4))

        ctk.CTkLabel(r_size, text="Largura da Assinatura:").pack(side="left")
        self.sig_width_var = tk.IntVar(value=150) # Padrão inicial: 150pt
        
        self.sig_width_lbl = ctk.CTkLabel(r_size, text="150 pt", width=55, font=ctk.CTkFont(weight="bold"))
        self.sig_width_lbl.pack(side="right", padx=4)

        def _on_slider_change(val):
            self.sig_width_lbl.configure(text=f"{int(val)} pt")
            self._sig_redraw_positions() # Atualiza o desenho no preview instantaneamente

        self.sig_width_slider = ctk.CTkSlider(
            r_size, from_=40, to=350, number_of_steps=310,
            variable=self.sig_width_var, command=_on_slider_change
        )
        self.sig_width_slider.pack(side="left", fill="x", expand=True, padx=8)

        # Linha 4: Página + Controles do Preview
        pg_r = ctk.CTkFrame(ctrl, fg_color="transparent")
        pg_r.pack(fill="x", pady=(8, 2))

        ctk.CTkLabel(pg_r, text="Página:").pack(side="left")
        self.sig_page_entry = ctk.CTkEntry(pg_r, width=60, placeholder_text="1")
        self.sig_page_entry.pack(side="left", padx=8)

        self._action_btn(
            pg_r, "🔍 Carregar Preview", self._sig_load_preview, width=140, color="#1E3A5F"
        ).pack(side="left", padx=4)

        self._action_btn(
            pg_r, "↩ Último ponto", self._sig_undo_point, width=110, color="#334155"
        ).pack(side="left", padx=4)

        self._action_btn(
            pg_r, "🗑 Limpar pontos", self._sig_clear_points, width=120, color="#7F1D1D"
        ).pack(side="left", padx=4)

        self.sig_pos_lbl = ctk.CTkLabel(
            ctrl, text="📍 Clique no preview para adicionar pontos de assinatura",
            text_color=TEXT_DIM, anchor="w"
        )
        self.sig_pos_lbl.pack(anchor="w", pady=(6, 0))

        # Preview vertical Expandido
        pv_frame = ctk.CTkFrame(cols, fg_color="#12121E", corner_radius=10)
        pv_frame.grid(row=0, column=1, sticky="nsew", padx=(4, 8))

        ctk.CTkLabel(
            pv_frame, text="Preview vertical — clique para marcar pontos",
            font=ctk.CTkFont("Segoe UI", 11), text_color=TEXT_DIM
        ).pack(pady=(8, 2))

        self.sig_preview_cv = tk.Canvas(
            pv_frame, bg="#12121E", cursor="crosshair", highlightthickness=0
        )
        self.sig_preview_cv.pack(fill="both", expand=True, padx=4, pady=4)
        self.sig_preview_cv.bind("<Button-1>", self._sig_on_click)

        # Transformação canvas ↔ PDF
        self._sig_preview_ratio = 1.0
        self._sig_preview_offx = 0
        self._sig_preview_offy = 0
        self._sig_preview_draw_w = 0
        self._sig_preview_draw_h = 0
        self._sig_pdf_w = 0
        self._sig_pdf_h = 0

        self._action_btn(
            ctrl, "✍️ Assinar e Salvar", self._sig_run, width=220
        ).pack(pady=14)


    def _sig_import_img(self):
        f = filedialog.askopenfilename(
            filetypes=[("Imagens", "*.png *.jpg *.jpeg")]
        )
        if f:
            self._sig_img_path = f
            self.sig_img_var.set(Path(f).name)
            try:
                # Carrega a imagem base em cache para usar em tempo real no canvas
                self._sig_pil_base = Image.open(f).convert("RGBA")
                self._sig_redraw_positions()
            except Exception as e:
                print(f"Erro ao carregar imagem para preview: {e}")


    def _sig_select_pdf(self):
        fs = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if fs:
            self._sig_pdf_paths = list(fs)
            self._sig_pdf_preview_path = self._sig_pdf_paths[0]
            self.sig_file_var.set(f"{len(fs)} PDF(s) selecionado(s)")
            self._sig_load_preview()


    def _sig_clear_points(self):
        self._sig_positions.clear()
        if hasattr(self, "sig_preview_cv"):
            self.sig_preview_cv.delete("sig_marker")
        self._preview_tkimgs.clear()
        self.sig_pos_lbl.configure(
            text="📍 Clique no preview para adicionar pontos de assinatura"
        )


    def _sig_undo_point(self):
        if self._sig_positions:
            self._sig_positions.pop()
            self._sig_redraw_positions()
            self._sig_update_pos_label()


    def _sig_update_pos_label(self):
        n = len(self._sig_positions)
        if n == 0:
            self.sig_pos_lbl.configure(
                text="📍 Clique no preview para adicionar pontos de assinatura"
            )
            return

        rx, ry = self._sig_positions[-1]
        x_pt = rx * self._sig_pdf_w if self._sig_pdf_w else 0
        y_pt = ry * self._sig_pdf_h if self._sig_pdf_h else 0
        self.sig_pos_lbl.configure(
            text=f"📍 {n} ponto(s) marcado(s) | último: x={x_pt:.0f} pt, y={y_pt:.0f} pt"
        )


    def _sig_redraw_positions(self):
        if not hasattr(self, "sig_preview_cv"):
            return

        self.sig_preview_cv.delete("sig_marker")
        self._preview_tkimgs.clear() # Limpa referências antigas para liberar memória

        if self._sig_preview_draw_w <= 0 or self._sig_preview_draw_h <= 0:
            return

        ox = self._sig_preview_offx
        oy = self._sig_preview_offy
        dw = self._sig_preview_draw_w
        dh = self._sig_preview_draw_h
        ratio = self._sig_preview_ratio 

        # Descobre a proporção da assinatura para manter o aspecto correto no desenho
        if self._sig_pil_base:
            sw, sh = self._sig_pil_base.size
            sig_ratio = sh / sw
        else:
            sig_ratio = 0.4 # Proporção retangular genérica padrão

        # Pega a largura atualizada do slider e converte proporcionalmente para o Canvas
        pdf_w_box = float(self.sig_width_var.get())
        pdf_h_box = pdf_w_box * sig_ratio

        cv_w_box = max(2, int(pdf_w_box * ratio))
        cv_h_box = max(2, int(pdf_h_box * ratio))

        for i, (rx, ry) in enumerate(self._sig_positions, 1):
            x = ox + (rx * dw)
            y = oy + (ry * dh)

            # MUDANÇA: Se já carregou a assinatura, renderiza ela em tempo real na tela!
            if self._sig_pil_base:
                try:
                    img_resized = self._sig_pil_base.resize((cv_w_box, cv_h_box), Image.LANCZOS)
                    tkimg = ImageTk.PhotoImage(img_resized)
                    self._preview_tkimgs.append(tkimg) # Trava referência
                    
                    self.sig_preview_cv.create_image(
                        x, y, image=tkimg, anchor="center", tags="sig_marker"
                    )
                except Exception:
                    pass
            else:
                # Se não houver imagem ainda, desenha a caixa de simulação pontilhada
                self.sig_preview_cv.create_rectangle(
                    x - cv_w_box/2, y - cv_h_box/2,
                    x + cv_w_box/2, y + cv_h_box/2,
                    outline="#EF4444", dash=(4, 4), width=2, tags="sig_marker"
                )

            # Indicadores visuais de centro e número da ordem
            self.sig_preview_cv.create_oval(
                x - 3, y - 3, x + 3, y + 3,
                fill="#EF4444", outline="#FFFFFF", width=1, tags="sig_marker"
            )
            self.sig_preview_cv.create_text(
                x + (cv_w_box/2) + 10, y - (cv_h_box/2) - 4,
                text=str(i), fill="#EF4444", font=("Segoe UI", 10, "bold"), tags="sig_marker"
            )


    def _sig_load_preview(self):
        pdf = self._sig_pdf_preview_path or (
            self._sig_pdf_paths[0] if self._sig_pdf_paths else None
        )
        if not pdf:
            messagebox.showwarning("Aviso", "Selecione PDFs primeiro.")
            return

        self._sig_pdf_preview_path = pdf

        try:
            page_num = int(self.sig_page_entry.get() or "1") - 1
        except ValueError:
            page_num = 0

        self.sig_preview_cv.update_idletasks()
        cw = self.sig_preview_cv.winfo_width()
        ch = self.sig_preview_cv.winfo_height()
        
        if cw <= 1 or ch <= 1:
            cw, ch = 600, 800

        def worker():
            try:
                with fitz.open(pdf) as doc:
                    if page_num < 0 or page_num >= len(doc):
                        raise ValueError(
                            f"O PDF {Path(pdf).name} não possui a página {page_num + 1}."
                        )
                    page = doc[page_num]
                    self._sig_pdf_w = page.rect.width
                    self._sig_pdf_h = page.rect.height

                    pix = page.get_pixmap(dpi=150)

                img = Image.open(BytesIO(pix.tobytes("png"))).convert("RGB")

                r = min(cw / img.width, ch / img.height)
                nw = max(1, int(img.width * r))
                nh = max(1, int(img.height * r))
                img = img.resize((nw, nh), Image.LANCZOS)

                self._sig_preview_draw_w = nw
                self._sig_preview_draw_h = nh
                self._sig_preview_ratio = nw / self._sig_pdf_w if self._sig_pdf_w else 1.0
                self._sig_preview_offx = (cw - nw) // 2
                self._sig_preview_offy = (ch - nh) // 2

                tkimg = ImageTk.PhotoImage(img)

                def ui():
                    self.sig_preview_cv.delete("all")
                    self.sig_preview_cv.create_image(
                        cw // 2, ch // 2,
                        image=tkimg,
                        anchor="center"
                    )
                    self.sig_preview_cv.image = tkimg
                    self._sig_clear_points()

                self.after(0, ui)

            except Exception as e:
                self.after(0, lambda: messagebox.showerror("Erro", str(e)))

        run_in_thread(worker)


    def _sig_on_click(self, ev):
        if self._sig_preview_draw_w <= 0 or self._sig_preview_draw_h <= 0:
            return

        ox = self._sig_preview_offx
        oy = self._sig_preview_offy
        dw = self._sig_preview_draw_w
        dh = self._sig_preview_draw_h

        if ev.x < ox or ev.y < oy or ev.x > (ox + dw) or ev.y > (oy + dh):
            return

        rx = (ev.x - ox) / dw
        ry = (ev.y - oy) / dh  

        self._sig_positions.append((rx, ry))
        self._sig_redraw_positions()
        self._sig_update_pos_label()


    def _sig_run(self):
        if not self._sig_pdf_paths:
            messagebox.showwarning("Aviso", "Selecione um ou mais PDFs.")
            return

        if not self._sig_positions:
            messagebox.showwarning("Aviso", "Marque ao menos um ponto no preview.")
            return

        if self._sig_img_path:
            sig_pil = Image.open(self._sig_img_path).convert("RGBA")
        else:
            messagebox.showwarning("Aviso", "Selecione uma imagem de assinatura.")
            return

        out_dir = filedialog.askdirectory(title="Selecione a pasta de saída")
        if not out_dir:
            return

        try:
            page_num = int(self.sig_page_entry.get() or "1") - 1
        except ValueError:
            page_num = 0

        def worker():
            self.set_status("Inserindo assinatura...", 0.3)
            try:
                tmp = Path(tempfile.gettempdir()) / "_cmat_sig.png"
                sig_pil.save(str(tmp), "PNG")

                total = len(self._sig_pdf_paths)
                ok = 0
                erros = 0

                sig_w, sig_h = sig_pil.size
                sig_ratio = (sig_h / sig_w) if sig_w else 0.4

                # MUDANÇA: O tamanho base final agora é o valor exato escolhido pelo usuário no slider
                base_w = float(self.sig_width_var.get())
                base_h = base_w * sig_ratio

                for i, pdf_path in enumerate(self._sig_pdf_paths, 1):
                    doc = None
                    try:
                        doc = fitz.open(pdf_path)

                        if page_num < 0 or page_num >= len(doc):
                            raise ValueError(
                                f"O PDF {Path(pdf_path).name} não possui a página {page_num + 1}."
                            )

                        page = doc[page_num]
                        pdf_w = page.rect.width
                        pdf_h = page.rect.height

                        for rx, ry in self._sig_positions:
                            cx = rx * pdf_w
                            cy = ry * pdf_h

                            w = min(base_w, pdf_w - 10)
                            h = min(base_h, pdf_h - 10)

                            x0 = cx - (w / 2)
                            y0 = cy - (h / 2)

                            x0 = max(0, min(x0, max(0, pdf_w - w)))
                            y0 = max(0, min(y0, max(0, pdf_h - h)))

                            rect = fitz.Rect(x0, y0, x0 + w, y0 + h)
                            page.insert_image(rect, filename=str(tmp), overlay=True)

                        out = unique_path(
                            Path(out_dir) / f"{Path(pdf_path).stem}_assinado.pdf"
                        )
                        doc.save(str(out))
                        ok += 1

                        self.set_status(
                            f"Assinando PDFs... ({i}/{total})",
                            i / total
                        )

                    except Exception as e:
                        erros += 1
                        print(f"Erro em {pdf_path}: {e}")
                    finally:
                        if doc is not None:
                            doc.close()

                self.after(0, lambda: messagebox.showinfo(
                    "Concluído",
                    f"{ok} PDF(s) assinados.\n{erros} com erro.\n\nSaída:\n{out_dir}"
                ))

            except Exception as e:
                self.after(0, lambda: messagebox.showerror("Erro", str(e)))

        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 🗜️ COMPRIMIR PDF
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_compress(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "🗜️", "Comprimir PDF",
                          "Reduza o tamanho de PDFs.")
        card = self._card(frame)
        card.pack(fill="both", expand=True, padx=28, pady=14)

        self.comp_files: List[str] = []
        self.comp_files_var = ctk.StringVar(value="Nenhum arquivo")
        r1 = ctk.CTkFrame(card, fg_color="transparent")
        r1.pack(fill="x", padx=16, pady=(14, 4))
        self._action_btn(r1, "📂 Selecionar PDFs",
                          self._comp_select, width=200).pack(side="left")
        ctk.CTkLabel(r1, textvariable=self.comp_files_var,
                     text_color=TEXT_DIM).pack(side="left", padx=10)

        ctk.CTkLabel(card, text="Nível de compressão:",
                     font=ctk.CTkFont("Segoe UI", 13, "bold"),
                     anchor="w").pack(anchor="w", padx=16, pady=(12, 4))
        self.comp_level = ctk.StringVar(value="media")
        lf = ctk.CTkFrame(card, fg_color="transparent")
        lf.pack(fill="x", padx=16)
        for val, lbl, desc in [
            ("leve",  "🟢 Leve",    "72 DPI / qualidade 80"),
            ("media", "🟡 Média",   "60 DPI / qualidade 60"),
            ("forte", "🔴 Intensa", "48 DPI / qualidade 40"),
        ]:
            row = ctk.CTkFrame(lf, fg_color="transparent")
            row.pack(fill="x", pady=2)
            ctk.CTkRadioButton(row, text=lbl, variable=self.comp_level,
                                value=val, width=120).pack(side="left")
            ctk.CTkLabel(row, text=desc, text_color=TEXT_DIM,
                         font=ctk.CTkFont("Segoe UI", 12)).pack(
                side="left", padx=8)

        self.comp_log = self._file_listbox(card, height=110)
        self.comp_log.pack(fill="both", expand=True, padx=16, pady=8)

        self._action_btn(card, "🗜️ Comprimir PDFs",
                          self._comp_run, width=200).pack(pady=(4, 14))
        self.comp_prog = ctk.CTkProgressBar(card)
        self.comp_prog.pack(fill="x", padx=16, pady=(0, 14))
        self.comp_prog.set(0)

    def _comp_select(self):
        fs = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if fs:
            self.comp_files = list(fs)
            self.comp_files_var.set(f"{len(fs)} arquivo(s)")

    def _comp_run(self):
        if not self.comp_files:
            messagebox.showwarning("Aviso", "Selecione PDFs."); return
        dpi, q = {"leve":(72,80),"media":(60,60),"forte":(48,40)
                  }[self.comp_level.get()]
        out_dir = filedialog.askdirectory(title="Pasta de saída")
        if not out_dir: return
        self.comp_log.delete("1.0", "end")

        def log(m):
            self.after(0, lambda msg=m: (
                self.comp_log.insert("end", msg+"\n"),
                self.comp_log.see("end")))

        def worker():
            total = len(self.comp_files)
            for i, fp in enumerate(self.comp_files):
                try:
                    om  = os.path.getsize(fp)/1_048_576
                    op  = unique_path(
                        Path(out_dir)/(Path(fp).stem+"_comprimido.pdf"))
                    red = compress_pdf_fitz(fp, str(op),
                                             image_dpi=dpi, jpeg_quality=q)
                    nm  = os.path.getsize(str(op))/1_048_576
                    log(f"{Path(fp).name}: {om:.2f}MB → {nm:.2f}MB "
                        f"(redução: {max(red,0)}%)")
                except Exception as e:
                    log(f"ERRO {Path(fp).name}: {e}")
                self.comp_prog.set((i+1)/total)
            self.set_status("Compressão concluída!", 1)
            self.after(0, lambda: messagebox.showinfo(
                "Concluído", f"PDFs comprimidos em:\n{out_dir}"))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 📦 DESCOMPRIMIR PDF  — nova função
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_decompress(self):
        self.clear_main()
        frame = self._page()
        self.current_frame = frame
        self._page_title(frame, "📦", "Descomprimir PDF",
                          "Remove a compressão interna de streams do PDF.")
        card = self._card(frame)
        card.pack(fill="both", expand=True, padx=28, pady=14)

        ctk.CTkLabel(card,
            text=(
                "ℹ️  Expande os streams comprimidos do PDF, tornando o arquivo "
                "maior porém mais acessível a editores externos, ferramentas "
                "de extração e análise forense.\n\n"
                "Use antes de editar PDFs em outros programas ou quando "
                "ferramentas de extração não consigam ler o conteúdo."
            ),
            wraplength=700, justify="left",
            font=ctk.CTkFont("Segoe UI", 12),
            text_color=TEXT_DIM).pack(anchor="w", padx=16, pady=(14, 8))

        self.decomp_files: List[str] = []
        self.decomp_files_var = ctk.StringVar(value="Nenhum arquivo")
        r1 = ctk.CTkFrame(card, fg_color="transparent")
        r1.pack(fill="x", padx=16, pady=4)
        self._action_btn(r1, "📂 Selecionar PDFs",
                          self._decomp_select, width=200).pack(side="left")
        ctk.CTkLabel(r1, textvariable=self.decomp_files_var,
                     text_color=TEXT_DIM).pack(side="left", padx=10)

        self.decomp_log = self._file_listbox(card, height=120)
        self.decomp_log.pack(fill="both", expand=True, padx=16, pady=8)

        self._action_btn(card, "📦 Descomprimir PDFs",
                          self._decomp_run, width=220).pack(pady=(4, 14))
        self.decomp_prog = ctk.CTkProgressBar(card)
        self.decomp_prog.pack(fill="x", padx=16, pady=(0, 14))
        self.decomp_prog.set(0)

    def _decomp_select(self):
        fs = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if fs:
            self.decomp_files = list(fs)
            self.decomp_files_var.set(f"{len(fs)} arquivo(s)")

    def _decomp_run(self):
        if not self.decomp_files:
            messagebox.showwarning("Aviso", "Selecione PDFs."); return
        out_dir = filedialog.askdirectory(title="Pasta de saída")
        if not out_dir: return
        self.decomp_log.delete("1.0", "end")

        def log(m):
            self.after(0, lambda msg=m: (
                self.decomp_log.insert("end", msg+"\n"),
                self.decomp_log.see("end")))

        def worker():
            total = len(self.decomp_files)
            for i, fp in enumerate(self.decomp_files):
                try:
                    op = unique_path(
                        Path(out_dir)/(Path(fp).stem+"_descomprimido.pdf"))
                    new_sz, orig_sz = decompress_pdf(fp, str(op))
                    fator = new_sz/orig_sz if orig_sz else 1
                    log(f"{Path(fp).name}: "
                        f"{orig_sz/1_048_576:.2f}MB → "
                        f"{new_sz/1_048_576:.2f}MB (×{fator:.1f})")
                except Exception as e:
                    log(f"ERRO {Path(fp).name}: {e}")
                self.decomp_prog.set((i+1)/total)
            self.set_status("Descompressão concluída!", 1)
            self.after(0, lambda: messagebox.showinfo(
                "Concluído", f"PDFs descomprimidos em:\n{out_dir}"))
        run_in_thread(worker)

    # ══════════════════════════════════════════════════════════════════════════
    # 🏷️ RENOMEADOR  — lógica original + Lista Auxiliar 2 + sem divisão de tela
    # ══════════════════════════════════════════════════════════════════════════

    def open_tab_renamer(self):
        self.clear_main()

        # Mesmo padrão do Home: frame externo → canvas scrollável
        outer = ctk.CTkFrame(self.main, fg_color=MAIN_BG)
        outer.pack(fill="both", expand=True)
        self.current_frame = outer          # salva UMA VEZ

        canvas  = tk.Canvas(outer, bg=MAIN_BG, highlightthickness=0)
        vscroll = tk.Scrollbar(outer, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vscroll.set)
        vscroll.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        inner  = ctk.CTkFrame(canvas, fg_color=MAIN_BG)
        win_id = canvas.create_window((0, 0), window=inner, anchor="nw")

        def _resize(e=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfig(win_id, width=canvas.winfo_width())

        inner.bind("<Configure>", _resize)
        canvas.bind("<Configure>", _resize)
        canvas.bind_all("<MouseWheel>",
            lambda e: canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        self.set_status("Renomeador de PDFs", 0)
        self._page_title(inner, "🏷️",
                          "Renomeador Inteligente de PDFs",
                          "Renomeie em lote usando OCR + correspondência fuzzy.")

        card = self._card(inner)
        card.pack(fill="x", padx=28, pady=14)

        ctk.CTkLabel(card, text="Renomeador de PDFs",
                     font=ctk.CTkFont(size=18, weight="bold")
                     ).pack(pady=(10, 10))

        # ── 3 listas lado a lado ─────────────────────────────────────────────
        lists_row = ctk.CTkFrame(card, fg_color="transparent")
        lists_row.pack(fill="x", padx=16, pady=(0, 10))
        lists_row.grid_columnconfigure((0, 1, 2), weight=1)

        ctk.CTkLabel(lists_row, text="Lista Principal:"
                     ).grid(row=0, column=0, sticky="w")
        self.main_list_text = ctk.CTkTextbox(lists_row, height=90, width=200)
        self.main_list_text.grid(row=1, column=0, padx=(0,8),
                                  pady=(0,8), sticky="ew")

        ctk.CTkLabel(lists_row, text="Lista Auxiliar 1 (opcional):"
                     ).grid(row=0, column=1, sticky="w")
        self.aux_list_text = ctk.CTkTextbox(lists_row, height=90, width=200)
        self.aux_list_text.grid(row=1, column=1, padx=(0,8),
                                 pady=(0,8), sticky="ew")

        ctk.CTkLabel(lists_row, text="Lista Auxiliar 2 (opcional):"
                     ).grid(row=0, column=2, sticky="w")
        self.aux_list2_text = ctk.CTkTextbox(lists_row, height=90, width=200)
        self.aux_list2_text.grid(row=1, column=2, pady=(0,8), sticky="ew")

        # ── Opções ───────────────────────────────────────────────────────────
        options_frame = ctk.CTkFrame(card)
        options_frame.pack(fill="x", padx=16, pady=(0, 10))

        ctk.CTkLabel(options_frame, text="Prefixo:"
                     ).grid(row=0, column=0, padx=5, pady=5)
        self.prefix_entry = ctk.CTkEntry(options_frame, width=120)
        self.prefix_entry.grid(row=0, column=1, padx=5, pady=5)

        ctk.CTkLabel(options_frame, text="Sufixo:"
                     ).grid(row=0, column=2, padx=5, pady=5)
        self.suffix_entry = ctk.CTkEntry(options_frame, width=120)
        self.suffix_entry.grid(row=0, column=3, padx=5, pady=5)

        ctk.CTkLabel(options_frame, text="Manter primeiros X:"
                     ).grid(row=1, column=0, padx=5, pady=5)
        self.keep_first_entry = ctk.CTkEntry(options_frame, width=80)
        self.keep_first_entry.insert(0, "0")
        self.keep_first_entry.grid(row=1, column=1, padx=5, pady=5)

        ctk.CTkLabel(options_frame, text="Manter últimos X:"
                     ).grid(row=1, column=2, padx=5, pady=5)
        self.keep_last_entry = ctk.CTkEntry(options_frame, width=80)
        self.keep_last_entry.insert(0, "0")
        self.keep_last_entry.grid(row=1, column=3, padx=5, pady=5)

        # ── Botões ────────────────────────────────────────────────────────────
        btn_frame = ctk.CTkFrame(card, fg_color="transparent")
        btn_frame.pack(pady=(4, 8))
        ctk.CTkButton(btn_frame, text="Selecionar PDFs",
                       command=self.select_pdfs_for_rename,
                       width=160).grid(row=0, column=0, padx=10, pady=5)
        ctk.CTkButton(btn_frame, text="Renomear Arquivos",
                       command=self.start_rename_process,
                       width=160).grid(row=0, column=1, padx=10, pady=5)

        # ── Preview ───────────────────────────────────────────────────────────
        ctk.CTkLabel(card, text="Pré-visualização dos novos nomes:",
                     anchor="w").pack(anchor="w", padx=16, pady=(4, 2))
        self.preview_box = ctk.CTkTextbox(card, height=160, width=550)
        self.preview_box.pack(fill="x", padx=16, pady=(0, 4))
        self.preview_box.insert("1.0", "Nenhum arquivo selecionado ainda.")

        self.rename_progress = ctk.CTkProgressBar(card)
        self.rename_progress.pack(fill="x", padx=16, pady=(4, 14))
        self.rename_progress.set(0)

        self.pdf_files_to_rename: List[str] = []

    # ── Lógica original preservada ────────────────────────────────────────────

    def select_pdfs_for_rename(self):
        files = filedialog.askopenfilenames(
            title="Selecionar PDFs",
            filetypes=[("PDF files", "*.pdf")])
        if not files: return
        self.pdf_files_to_rename = list(files)
        self.preview_box.delete("1.0", "end")
        self.preview_box.insert("1.0", "Gerando pré-visualização...\n")
        threading.Thread(target=self.gen_preview, daemon=True).start()

    def gen_preview(self):
        """Gera pré-visualização dos nomes — lógica original + lista auxiliar 2."""
        main_list  = [x.strip() for x in
                      self.main_list_text.get("1.0","end").splitlines()
                      if x.strip()]
        aux_list   = [x.strip() for x in
                      self.aux_list_text.get("1.0","end").splitlines()
                      if x.strip()]
        aux_list2  = [x.strip() for x in
                      self.aux_list2_text.get("1.0","end").splitlines()
                      if x.strip()]
        prefix = self.prefix_entry.get().strip()
        suffix = self.suffix_entry.get().strip()
        try:    keep_first = int(self.keep_first_entry.get() or 0)
        except: keep_first = 0
        try:    keep_last  = int(self.keep_last_entry.get()  or 0)
        except: keep_last  = 0

        files = list(self.pdf_files_to_rename)
        total = len(files)
        if total == 0: return

        self.rename_progress.set(0)
        self.preview_box.delete("1.0", "end")
        self.preview_box.insert("1.0",
            "Gerando pré-visualização (lendo arquivos)...\n")

        preview_lines = []

        for i, file in enumerate(files, 1):
            try:
                txt = ""
                with fitz.open(file) as doc:
                    for page in doc:
                        txt += page.get_text("text") or ""
                        # OCR se pouco texto
                        if len(txt.strip()) < 30 and HAVE_TESSERACT:
                            pix = page.get_pixmap(dpi=200)
                            img = Image.frombytes("RGB",
                                [pix.width, pix.height], pix.samples)
                            try:
                                txt += pytesseract.image_to_string(
                                    img, lang="por", config="--psm 6")
                            except Exception:
                                pass

                match_main  = self._find_best_match(txt, main_list)
                match_aux   = self._find_best_match(txt, aux_list)
                match_aux2  = self._find_best_match(txt, aux_list2)
                orig_stem   = Path(file).stem

                parts = []
                if prefix:          parts.append(prefix)
                if keep_first > 0:  parts.append(orig_stem[:keep_first])
                if match_main:      parts.append(match_main)
                if match_aux:       parts.append(match_aux)
                if match_aux2:      parts.append(match_aux2)
                if keep_last > 0:   parts.append(orig_stem[-keep_last:])
                if suffix:          parts.append(suffix)

                parts = [p for p in parts if p and str(p).strip()]
                if parts:
                    new_name = "_".join(parts)
                else:
                    if prefix or suffix:
                        mid = orig_stem if (keep_first>0 or keep_last>0) else ""
                        new_name = "_".join(
                            [p for p in [prefix, mid, suffix] if p])
                    else:
                        new_name = orig_stem

                new_name = re.sub(r'\s+', '_', new_name).strip("_")
                new_name = f"{new_name}.pdf"
                preview_lines.append((file, new_name))
                self.preview_box.insert(
                    "end", f"{Path(file).name}  →  {new_name}\n")

            except Exception as e:
                self.preview_box.insert(
                    "end", f"Erro ao ler {Path(file).name}: {e}\n")

            self.rename_progress.set(i / total * 0.5)

        self._preview_lines = preview_lines
        self.rename_progress.set(0.5)
        self.preview_box.insert(
            "end",
            "\nPré-visualização gerada. "
            "Clique em 'Renomear Arquivos' para aplicar.\n")

    def _find_best_match(self, text: str, candidates: List[str]):
        """Correspondência fuzzy — lógica original."""
        best, best_score = None, 0.0
        for c in candidates:
            if not c: continue
            if USE_RAPIDFUZZ:
                score = rfuzz.partial_ratio(c.lower(), text.lower())
            else:
                score = difflib.SequenceMatcher(
                    None, c.lower(), text.lower()).ratio() * 100
            if score > best_score and score >= 70:
                best, best_score = c, score
        return best

    def start_rename_process(self):
        if not getattr(self, "_preview_lines", None):
            messagebox.showwarning("Aviso",
                "Gere a pré-visualização antes de renomear "
                "(clique em 'Selecionar PDFs').")
            return
        outdir = filedialog.askdirectory(title="Pasta de saída")
        if not outdir: return
        threading.Thread(target=self._rename_files,
                          args=(outdir,), daemon=True).start()

    def _rename_files(self, outdir: str):
        lines = getattr(self, "_preview_lines", [])
        total = len(lines)
        if total == 0: return
        outdir = Path(outdir)
        outdir.mkdir(parents=True, exist_ok=True)
        self.rename_progress.set(0.5)
        renamed = 0
        for i, (src, new_name) in enumerate(lines, 1):
            try:
                base = Path(new_name).stem
                dest = outdir / f"{base}.pdf"
                counter = 1
                while dest.exists():
                    dest = outdir / f"{base}({counter}).pdf"
                    counter += 1
                shutil.copy2(src, dest)
                renamed += 1
            except Exception as e:
                print("rename err:", src, e)
            self.rename_progress.set(0.5 + (i/total)*0.5)
        msg = (f"Operação finalizada: {renamed} arquivo(s) "
               f"processado(s).\nDestino: {outdir}")
        self.after(0, lambda: messagebox.showinfo("Concluído", msg))
        self.rename_progress.set(1.0)
        self.set_status(msg, 1)


# ── Entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    app = App()
    app.mainloop()